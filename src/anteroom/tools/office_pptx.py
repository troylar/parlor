"""PPTX (PowerPoint presentation) create/read/edit tool.

Backends:
- COM (Windows + Office + pywin32): full Office object model
- Library (python-pptx): cross-platform XML manipulation

Install: ``pip install anteroom[office]`` or ``pip install anteroom[office-com]``
"""

from __future__ import annotations

import os
import sys
from typing import Any, Optional

from .security import validate_path

_BACKEND: str | None = None
_com_mod: Any = None

if sys.platform == "win32":
    try:
        from . import office_com as _com_mod_import

        if _com_mod_import.COM_AVAILABLE:
            _com_mod = _com_mod_import
            _BACKEND = "com"
    except ImportError:
        pass

if _BACKEND is None:
    try:
        from pptx import Presentation  # noqa: F401

        _BACKEND = "lib"
    except ImportError:
        pass

AVAILABLE = _BACKEND is not None

_MAX_OUTPUT = 100_000
_MAX_CONTENT_BLOCKS = 200
_MAX_SLIDES = 100
_MAX_EDIT_OPS = 500

# English Metric Units per inch (used by COM backend for position/size)
_EMU_PER_INCH = 914400

# Shared shape type names supported by both backends
_SHAPE_TYPE_NAMES = (
    "rectangle",
    "oval",
    "triangle",
    "right_arrow",
    "left_arrow",
    "diamond",
    "rounded_rectangle",
    "star",
)

# COM MsoAutoShapeType constants (shared between insert_shape and other COM callers)
_COM_SHAPE_TYPE_MAP: dict[str, int] = {
    "rectangle": 1,
    "oval": 9,
    "triangle": 7,
    "right_arrow": 33,
    "left_arrow": 34,
    "diamond": 4,
    "rounded_rectangle": 5,
    "star": 12,
}

_ALL_ACTIONS = [
    "create",
    "read",
    "edit",
    "transitions",
    "animations",
    "insert_image",
    "insert_shape",
    "format_shape",
    "master_layout",
    "reorder_slides",
    "embed_chart",
    "embed_table",
    "export_pdf",
    "hyperlinks",
    "headers_footers",
    "sections",
    "group_shapes",
    "audio_video",
    "smartart",
]

_working_dir: str = os.getcwd()

DEFINITION: dict[str, Any] = {
    "name": "pptx",
    "description": (
        "Create, read, edit, and manipulate PowerPoint presentations (.pptx). "
        "Supports slide transitions, animations, images, shapes, formatting, charts, "
        "tables, PDF export, hyperlinks, headers/footers, sections, and more."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": _ALL_ACTIONS,
                "description": "Operation to perform.",
            },
            "path": {
                "type": "string",
                "description": "File path (relative to working directory or absolute).",
            },
            "slides": {
                "type": "array",
                "description": (
                    "Slides for create/edit append. Each slide: "
                    "{title?: str, content?: str, bullets?: [str], notes?: str, layout?: int}. "
                    "layout: 0=title slide, 1=title+content (default), 5=blank, 6=blank."
                ),
                "items": {"type": "object"},
            },
            "replacements": {
                "type": "array",
                "description": "List of {old: str, new: str} for find/replace in edit action.",
                "items": {"type": "object"},
            },
            "table_edits": {
                "type": "array",
                "description": (
                    "Edit cells in existing tables. Each entry: "
                    "{slide_index: int (1-based), table_index?: int (1-based, default 1), "
                    "row: int (0-based), col: int (0-based), value: str}. "
                    "Use action='read' first to discover table positions."
                ),
                "items": {"type": "object"},
            },
            "shape_edits": {
                "type": "array",
                "description": (
                    "Edit text/formatting on specific shapes. Each entry: "
                    "{slide_index: int (1-based), shape_index: int (1-based), "
                    "text?: str, font_size?: number, font_bold?: bool, "
                    "font_color?: str (hex), font_name?: str}."
                ),
                "items": {"type": "object"},
            },
            "notes_edits": {
                "type": "array",
                "description": (
                    "Set speaker notes on specific slides. Each entry: {slide_index: int (1-based), text: str}."
                ),
                "items": {"type": "object"},
            },
            "delete_slides": {
                "type": "array",
                "description": "List of 1-based slide indices to delete. Processed in descending order.",
                "items": {"type": "integer"},
            },
            "duplicate_slides": {
                "type": "array",
                "description": (
                    "List of 1-based slide indices to duplicate. "
                    "Each duplicate is inserted immediately after the original."
                ),
                "items": {"type": "integer"},
            },
            "template_fill": {
                "type": "object",
                "description": (
                    "Key-value pairs for {{key}} token replacement throughout the entire presentation. "
                    "Replaces {{key}} with value in all text frames, tables, and notes. "
                    'Example: {"company": "Acme", "date": "2025-01-01"} replaces '
                    "{{company}} and {{date}} everywhere."
                ),
            },
            "table_format": {
                "type": "array",
                "description": (
                    "Format cells in existing tables. Each entry: "
                    "{slide_index: int (1-based), table_index?: int (1-based, default 1), "
                    "row: int (0-based), col: int (0-based), "
                    "bg_color?: str (hex), font_size?: number, font_bold?: bool, "
                    "font_color?: str (hex), font_name?: str, "
                    "alignment?: str ('left'|'center'|'right')}."
                ),
                "items": {"type": "object"},
            },
            "paragraph_edits": {
                "type": "array",
                "description": (
                    "Set multi-paragraph text on a shape with per-paragraph formatting. Each entry: "
                    "{slide_index: int (1-based), shape_index: int (1-based), "
                    "paragraphs: [{text: str, level?: int (0-8), font_size?: number, "
                    "font_bold?: bool, font_color?: str (hex), font_name?: str, "
                    "alignment?: str ('left'|'center'|'right')}]}."
                ),
                "items": {"type": "object"},
            },
            "placeholder_edits": {
                "type": "array",
                "description": (
                    "Edit placeholders by type name instead of index. Each entry: "
                    "{slide_index: int (1-based), placeholder_type: str "
                    "('title'|'body'|'subtitle'|'center_title'|'slide_number'|'date'|'footer'), "
                    "text?: str, font_size?: number, font_bold?: bool, "
                    "font_color?: str (hex), font_name?: str}."
                ),
                "items": {"type": "object"},
            },
            "image_replacements": {
                "type": "array",
                "description": (
                    "Replace image shapes with new images, preserving position and size. Each entry: "
                    "{slide_index: int (1-based), shape_index: int (1-based), "
                    "image_path: str}. The shape at shape_index must be a picture."
                ),
                "items": {"type": "object"},
            },
            "slide_index": {
                "type": "integer",
                "description": "1-based slide index for actions targeting a specific slide.",
            },
            "transition": {
                "type": "object",
                "description": (
                    "Transition settings: {effect?: str, advance_on_time?: bool, "
                    "advance_time?: number (seconds), speed?: str ('fast'|'medium'|'slow')}."
                ),
            },
            "shape_index": {
                "type": "integer",
                "description": "1-based shape index on the target slide.",
            },
            "effect_id": {
                "type": "integer",
                "description": "Animation effect ID for animations action.",
            },
            "operation": {
                "type": "string",
                "enum": ["add", "list", "delete", "duplicate", "move", "set", "apply"],
                "description": "Sub-operation for actions that support multiple modes.",
            },
            "image_path": {
                "type": "string",
                "description": "Image file path for insert_image action.",
            },
            "left": {
                "type": "number",
                "description": "Left position in inches.",
            },
            "top": {
                "type": "number",
                "description": "Top position in inches.",
            },
            "width": {
                "type": "number",
                "description": "Width in inches.",
            },
            "height": {
                "type": "number",
                "description": "Height in inches.",
            },
            "shape_type": {
                "type": "string",
                "description": (
                    "Shape type for insert_shape: 'rectangle', 'oval', 'triangle', "
                    "'right_arrow', 'left_arrow', 'diamond', 'rounded_rectangle', 'star'."
                ),
            },
            "format": {
                "type": "object",
                "description": (
                    "Shape formatting for format_shape: {fill_color?: str (hex), "
                    "line_color?: str (hex), line_width?: number (pt), "
                    "text?: str, font_size?: number, font_bold?: bool, "
                    "font_color?: str (hex), shadow?: bool}."
                ),
            },
            "layout_index": {
                "type": "integer",
                "description": "Layout index for master_layout apply operation.",
            },
            "new_position": {
                "type": "integer",
                "description": "1-based destination position for reorder_slides move.",
            },
            "chart_type": {
                "type": "string",
                "description": "Chart type for embed_chart: 'bar', 'line', 'pie', 'column'.",
            },
            "chart_title": {
                "type": "string",
                "description": "Chart title text.",
            },
            "data": {
                "type": "array",
                "description": (
                    "Data for embed_chart or embed_table. "
                    "For tables: [[header1, header2, ...], [row1val1, row1val2, ...], ...]. "
                    "For charts: [{name: str, values: [num, ...]}] series list."
                ),
                "items": {"type": "object"},
            },
            "rows": {
                "type": "integer",
                "description": "Number of rows for embed_table.",
            },
            "cols": {
                "type": "integer",
                "description": "Number of columns for embed_table.",
            },
            "output_path": {
                "type": "string",
                "description": "Output file path for export_pdf.",
            },
            "url": {
                "type": "string",
                "description": "URL for hyperlinks action.",
            },
            "display_text": {
                "type": "string",
                "description": "Display text for hyperlinks.",
            },
            "footer_text": {
                "type": "string",
                "description": "Footer text for headers_footers action.",
            },
            "slide_numbers": {
                "type": "boolean",
                "description": "Enable slide numbers (headers_footers action).",
            },
            "date_time": {
                "type": "string",
                "description": "Date/time text for headers_footers action.",
            },
            "section_name": {
                "type": "string",
                "description": "Section name for sections action.",
            },
            "shape_indices": {
                "type": "array",
                "description": "List of 1-based shape indices for group_shapes action.",
                "items": {"type": "integer"},
            },
            "media_path": {
                "type": "string",
                "description": "Media file path for audio_video action.",
            },
            "smartart_layout": {
                "type": "string",
                "description": "SmartArt layout name for smartart action.",
            },
            "smartart_items": {
                "type": "array",
                "description": "Text items for SmartArt nodes.",
                "items": {"type": "string"},
            },
        },
        "required": ["action", "path"],
    },
}


def set_working_dir(d: str) -> None:
    global _working_dir
    _working_dir = d


def _get_working_dir() -> str:
    """Return the current working directory, preferring the explicitly-set value."""
    return _working_dir


def _com_only_error(action: str) -> dict[str, Any]:
    """Return a descriptive error for actions that require the COM backend."""
    return {"error": f"Action '{action}' requires Windows with Office installed (COM backend)"}


async def handle(action: str, path: str, **kwargs: Any) -> dict[str, Any]:
    if not AVAILABLE:
        return {"error": "No pptx backend available. Install with: pip install anteroom[office]"}

    working_dir = _get_working_dir()
    resolved, error = validate_path(path, working_dir)
    if error:
        return {"error": error}

    if _BACKEND == "com":
        return await _dispatch_com(action, resolved, path, working_dir=working_dir, **kwargs)

    # Library backend dispatch
    _lib_dispatch: dict[str, Any] = {
        "create": _create_lib,
        "read": _read_lib,
        "edit": _edit_lib,
        "transitions": _transitions_lib,
        "animations": _animations_lib,
        "insert_image": _insert_image_lib,
        "insert_shape": _insert_shape_lib,
        "format_shape": _format_shape_lib,
        "master_layout": _master_layout_lib,
        "reorder_slides": _reorder_slides_lib,
        "embed_chart": _embed_chart_lib,
        "embed_table": _embed_table_lib,
        "export_pdf": _export_pdf_lib,
        "hyperlinks": _hyperlinks_lib,
        "headers_footers": _headers_footers_lib,
        "sections": _sections_lib,
        "group_shapes": _group_shapes_lib,
        "audio_video": _audio_video_lib,
        "smartart": _smartart_lib,
    }

    handler = _lib_dispatch.get(action)
    if handler is None:
        return {"error": f"Unknown action: {action}. Available: {', '.join(_ALL_ACTIONS)}"}
    return handler(resolved, path, working_dir=working_dir, **kwargs)


# ---------------------------------------------------------------------------
# COM backend
# ---------------------------------------------------------------------------


async def _dispatch_com(
    action: str, resolved: str, display_path: str, *, working_dir: str, **kwargs: Any
) -> dict[str, Any]:
    manager = _com_mod.get_manager()

    _com_dispatch: dict[str, Any] = {
        "create": _create_com,
        "read": _read_com,
        "edit": _edit_com,
        "transitions": _transitions_com,
        "animations": _animations_com,
        "insert_image": _insert_image_com,
        "insert_shape": _insert_shape_com,
        "format_shape": _format_shape_com,
        "master_layout": _master_layout_com,
        "reorder_slides": _reorder_slides_com,
        "embed_chart": _embed_chart_com,
        "embed_table": _embed_table_com,
        "export_pdf": _export_pdf_com,
        "hyperlinks": _hyperlinks_com,
        "headers_footers": _headers_footers_com,
        "sections": _sections_com,
        "group_shapes": _group_shapes_com,
        "audio_video": _audio_video_com,
        "smartart": _smartart_com,
    }

    handler = _com_dispatch.get(action)
    if handler is None:
        return {"error": f"Unknown action: {action}. Available: {', '.join(_ALL_ACTIONS)}"}
    try:
        return await manager.run_com(handler, manager, resolved, display_path, **kwargs)
    except Exception as exc:
        return {"error": f"COM {action} failed on {display_path}: {type(exc).__name__}: {exc}"}


def _open_pres_com(
    manager: Any,
    resolved: str,
    display_path: str,
    read_only: bool = False,
) -> tuple[Any, Any, Optional[str]]:
    """Open presentation via COM. Returns (ppt, prs, error).

    Callers MUST check the error string before using ppt/prs — both will
    be ``None`` when an error is returned.
    """
    if not os.path.isfile(resolved):
        return None, None, f"File not found: {display_path}"
    ppt = manager.get_app("PowerPoint.Application")
    try:
        prs = ppt.Presentations.Open(
            os.path.abspath(resolved),
            ReadOnly=read_only,
            WithWindow=True,
        )
        return ppt, prs, None
    except Exception as exc:
        return None, None, f"Unable to open PPTX file: {display_path} ({type(exc).__name__}: {exc})"


def _get_slide_com(prs: Any, slide_index: Optional[int]) -> tuple[Any, Optional[str]]:
    """Get a slide by 1-based index. Returns (slide, error)."""
    if slide_index is None:
        return None, "slide_index is required"
    if slide_index < 1 or slide_index > prs.Slides.Count:
        return None, f"slide_index {slide_index} out of range (1-{prs.Slides.Count})"
    return prs.Slides(slide_index), None


# --- create/read/edit (original) ---


_PP_LAYOUT_TEXT = 2  # ppLayoutText constant for COM


def _add_slide_com(prs: Any, slide_def: dict[str, Any]) -> None:
    """Add a single slide to a COM presentation."""
    layout_idx = slide_def.get("layout", _PP_LAYOUT_TEXT)
    # COM layout indices: 1=Title, 2=Title+Content, 7=Blank
    try:
        layout = prs.SlideMaster.CustomLayouts(layout_idx)
    except Exception:
        layout = prs.SlideMaster.CustomLayouts(2)

    slide = prs.Slides.AddSlide(prs.Slides.Count + 1, layout)

    title_text = slide_def.get("title")
    if title_text and slide.Shapes.HasTitle:
        slide.Shapes.Title.TextFrame.TextRange.Text = str(title_text)

    content = slide_def.get("content")
    bullets = slide_def.get("bullets")

    if content or bullets:
        # Find the body placeholder (index 2 in COM)
        body = None
        for i in range(1, slide.Shapes.Count + 1):
            shape = slide.Shapes(i)
            try:
                if shape.HasTextFrame and shape.PlaceholderFormat is not None and shape.PlaceholderFormat.Type == 2:
                    body = shape
                    break
            except Exception:
                continue

        if body is not None:
            tf = body.TextFrame.TextRange
            if content:
                tf.Text = str(content)
            elif bullets:
                tf.Text = str(bullets[0])
                for bullet in bullets[1:]:
                    tf.InsertAfter("\r" + str(bullet))

    notes_text = slide_def.get("notes")
    if notes_text:
        slide.NotesPage.Shapes(2).TextFrame.TextRange.Text = str(notes_text)


def _create_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    slides: list[dict[str, Any]] = kwargs.get("slides") or []
    if not slides:
        return {"error": "slides is required for create action"}
    if len(slides) > _MAX_SLIDES:
        return {"error": f"Too many slides (max {_MAX_SLIDES})"}

    ppt = manager.get_app("PowerPoint.Application")
    prs = ppt.Presentations.Add(WithWindow=True)

    try:
        for slide_def in slides:
            _add_slide_com(prs, slide_def)

        os.makedirs(os.path.dirname(resolved) or ".", exist_ok=True)
        # 24 = ppSaveAsOpenXMLPresentation (.pptx)
        prs.SaveAs(os.path.abspath(resolved), FileFormat=24)
    finally:
        prs.Close()

    return {"result": f"Created {display_path}", "path": display_path, "slides_created": len(slides)}


def _read_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path, read_only=True)
    if err:
        return {"error": err}

    try:
        output_parts: list[str] = []
        slide_count = prs.Slides.Count

        for i in range(1, slide_count + 1):
            slide = prs.Slides(i)
            output_parts.append(f"--- Slide {i} ({slide.Shapes.Count} shapes) ---")
            table_idx = 0
            for j in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(j)
                if shape.HasTable:
                    table_idx += 1
                    tbl = shape.Table
                    output_parts.append(f"[Table {table_idx}: {tbl.Rows.Count} rows x {tbl.Columns.Count} cols]")
                    for r in range(1, min(tbl.Rows.Count + 1, 51)):
                        row_vals: list[str] = []
                        for c in range(1, tbl.Columns.Count + 1):
                            cell_text = tbl.Cell(r, c).Shape.TextFrame.TextRange.Text.strip()
                            row_vals.append(cell_text)
                        output_parts.append(f"  Row {r - 1}: {row_vals}")
                elif shape.HasTextFrame:
                    text = shape.TextFrame.TextRange.Text.strip()
                    if text:
                        output_parts.append(f"[Shape {j}] {text}")
            if slide.HasNotesPage:
                try:
                    notes = slide.NotesPage.Shapes(2).TextFrame.TextRange.Text.strip()
                    if notes:
                        output_parts.append(f"[Notes] {notes}")
                except Exception:
                    pass

        content = "\n".join(output_parts)
        if len(content) > _MAX_OUTPUT:
            content = content[:_MAX_OUTPUT] + "\n... (truncated)"
    finally:
        prs.Close()

    return {
        "content": content,
        "slides": slide_count,
    }


def _edit_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    replacements: list[dict[str, str]] = kwargs.get("replacements") or []
    slides: list[dict[str, Any]] = kwargs.get("slides") or []
    table_edits: list[dict[str, Any]] = kwargs.get("table_edits") or []
    shape_edits: list[dict[str, Any]] = kwargs.get("shape_edits") or []
    notes_edits: list[dict[str, Any]] = kwargs.get("notes_edits") or []
    delete_slides: list[int] = kwargs.get("delete_slides") or []
    duplicate_slides: list[int] = kwargs.get("duplicate_slides") or []
    template_fill: dict[str, str] = kwargs.get("template_fill") or {}
    table_format: list[dict[str, Any]] = kwargs.get("table_format") or []
    paragraph_edits: list[dict[str, Any]] = kwargs.get("paragraph_edits") or []
    placeholder_edits: list[dict[str, Any]] = kwargs.get("placeholder_edits") or []
    image_replacements: list[dict[str, Any]] = kwargs.get("image_replacements") or []

    has_work = (
        replacements
        or slides
        or table_edits
        or shape_edits
        or notes_edits
        or delete_slides
        or duplicate_slides
        or template_fill
        or table_format
        or paragraph_edits
        or placeholder_edits
        or image_replacements
    )
    if not has_work:
        prs.Close()
        return {
            "error": (
                "Provide at least one of: 'replacements', 'slides', 'table_edits', "
                "'shape_edits', 'notes_edits', 'delete_slides', 'duplicate_slides', "
                "'template_fill', 'table_format', 'paragraph_edits', "
                "'placeholder_edits', 'image_replacements'"
            )
        }

    if slides and len(slides) > _MAX_SLIDES:
        prs.Close()
        return {"error": f"Too many slides to append (max {_MAX_SLIDES})"}

    for _arr_name, _arr in [
        ("table_edits", table_edits),
        ("table_format", table_format),
        ("shape_edits", shape_edits),
        ("notes_edits", notes_edits),
        ("paragraph_edits", paragraph_edits),
        ("placeholder_edits", placeholder_edits),
        ("image_replacements", image_replacements),
        ("delete_slides", delete_slides),
        ("duplicate_slides", duplicate_slides),
    ]:
        if len(_arr) > _MAX_EDIT_OPS:
            prs.Close()
            return {"error": f"Too many {_arr_name} entries (max {_MAX_EDIT_OPS})"}

    if len(template_fill) > _MAX_EDIT_OPS:
        prs.Close()
        return {"error": f"Too many template_fill keys (max {_MAX_EDIT_OPS})"}

    current_count = prs.Slides.Count
    if slides and current_count + len(slides) > _MAX_SLIDES:
        prs.Close()
        return {"error": f"Total slides would exceed limit (max {_MAX_SLIDES})"}

    try:
        result: dict[str, Any] = {"result": f"Edited {display_path}", "path": display_path}

        # Template fill — {{key}} replacement throughout entire presentation
        tokens_replaced = 0
        if template_fill:
            for key, value in template_fill.items():
                token = "{{" + str(key) + "}}"
                val_str = str(value)
                for i in range(1, prs.Slides.Count + 1):
                    slide = prs.Slides(i)
                    for j in range(1, slide.Shapes.Count + 1):
                        shape = slide.Shapes(j)
                        if shape.HasTextFrame:
                            tr = shape.TextFrame.TextRange
                            find = tr.Find(token)
                            while find is not None:
                                find.Text = val_str
                                tokens_replaced += 1
                                find = tr.Find(token)
                        if shape.HasTable:
                            tbl = shape.Table
                            for r in range(1, tbl.Rows.Count + 1):
                                for c in range(1, tbl.Columns.Count + 1):
                                    cell_tr = tbl.Cell(r, c).Shape.TextFrame.TextRange
                                    find = cell_tr.Find(token)
                                    while find is not None:
                                        find.Text = val_str
                                        tokens_replaced += 1
                                        find = cell_tr.Find(token)
                    if slide.HasNotesPage:
                        try:
                            notes_tr = slide.NotesPage.Shapes(2).TextFrame.TextRange
                            find = notes_tr.Find(token)
                            while find is not None:
                                find.Text = val_str
                                tokens_replaced += 1
                                find = notes_tr.Find(token)
                        except Exception:
                            pass
            result["tokens_replaced"] = tokens_replaced

        # Find/replace
        replacements_made = 0
        for rep in replacements:
            old = rep.get("old", "")
            new = rep.get("new", "")
            if not old:
                continue
            for i in range(1, prs.Slides.Count + 1):
                slide = prs.Slides(i)
                for j in range(1, slide.Shapes.Count + 1):
                    shape = slide.Shapes(j)
                    if shape.HasTextFrame:
                        tr = shape.TextFrame.TextRange
                        find = tr.Find(old)
                        while find is not None:
                            find.Text = new
                            replacements_made += 1
                            find = tr.Find(old)
        if replacements:
            result["replacements_made"] = replacements_made

        # Table cell edits
        table_cells_edited = 0
        for te in table_edits:
            si = te.get("slide_index")
            if si is None or si < 1 or si > prs.Slides.Count:
                continue
            slide = prs.Slides(si)
            ti = te.get("table_index", 1)
            row = te.get("row")
            col = te.get("col")
            value = te.get("value", "")
            if row is None or col is None:
                continue
            table_count = 0
            for j in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(j)
                if shape.HasTable:
                    table_count += 1
                    if table_count == ti:
                        tbl = shape.Table
                        if 0 <= row < tbl.Rows.Count and 0 <= col < tbl.Columns.Count:
                            tbl.Cell(row + 1, col + 1).Shape.TextFrame.TextRange.Text = str(value)
                            table_cells_edited += 1
                        break
        if table_edits:
            result["table_cells_edited"] = table_cells_edited

        # Table cell formatting
        table_cells_formatted = 0
        for tf_entry in table_format:
            si = tf_entry.get("slide_index")
            if si is None or si < 1 or si > prs.Slides.Count:
                continue
            slide = prs.Slides(si)
            ti = tf_entry.get("table_index", 1)
            row = tf_entry.get("row")
            col = tf_entry.get("col")
            if row is None or col is None:
                continue
            table_count = 0
            for j in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(j)
                if shape.HasTable:
                    table_count += 1
                    if table_count == ti:
                        tbl = shape.Table
                        if 0 <= row < tbl.Rows.Count and 0 <= col < tbl.Columns.Count:
                            cell = tbl.Cell(row + 1, col + 1)
                            bg_color = tf_entry.get("bg_color")
                            if bg_color:
                                cell.Shape.Fill.Visible = True
                                cell.Shape.Fill.ForeColor.RGB = _parse_color_int(bg_color)
                            alignment = tf_entry.get("alignment")
                            if alignment:
                                align_map = {"left": 1, "center": 2, "right": 3}
                                align_val = align_map.get(alignment.lower())
                                if align_val is not None:
                                    cell.Shape.TextFrame.TextRange.ParagraphFormat.Alignment = align_val
                            font_size = tf_entry.get("font_size")
                            if font_size is not None:
                                cell.Shape.TextFrame.TextRange.Font.Size = font_size
                            font_bold = tf_entry.get("font_bold")
                            if font_bold is not None:
                                cell.Shape.TextFrame.TextRange.Font.Bold = bool(font_bold)
                            font_color = tf_entry.get("font_color")
                            if font_color:
                                cell.Shape.TextFrame.TextRange.Font.Color.RGB = _parse_color_int(font_color)
                            font_name = tf_entry.get("font_name")
                            if font_name is not None:
                                cell.Shape.TextFrame.TextRange.Font.Name = str(font_name)
                            table_cells_formatted += 1
                        break
        if table_format:
            result["table_cells_formatted"] = table_cells_formatted

        # Paragraph edits — multi-paragraph text with per-paragraph formatting
        paragraphs_edited = 0
        for pe in paragraph_edits:
            si = pe.get("slide_index")
            shi = pe.get("shape_index")
            paras = pe.get("paragraphs")
            if si is None or shi is None or not paras:
                continue
            if si < 1 or si > prs.Slides.Count:
                continue
            slide = prs.Slides(si)
            if shi < 1 or shi > slide.Shapes.Count:
                continue
            shape = slide.Shapes(shi)
            if not shape.HasTextFrame:
                continue
            shape.TextFrame.TextRange.Text = ""
            for p_idx, p_def in enumerate(paras):
                text = p_def.get("text", "")
                tr = shape.TextFrame.TextRange
                if p_idx == 0:
                    tr.Text = str(text)
                else:
                    tr.InsertAfter("\r" + str(text))
                para_range = shape.TextFrame.TextRange.Paragraphs(p_idx + 1)
                level = p_def.get("level")
                if level is not None:
                    para_range.IndentLevel = max(1, min(9, int(level) + 1))
                alignment = p_def.get("alignment")
                if alignment:
                    align_map = {"left": 1, "center": 2, "right": 3}
                    align_val = align_map.get(alignment.lower())
                    if align_val is not None:
                        para_range.ParagraphFormat.Alignment = align_val
                font_size = p_def.get("font_size")
                if font_size is not None:
                    para_range.Font.Size = font_size
                font_bold = p_def.get("font_bold")
                if font_bold is not None:
                    para_range.Font.Bold = bool(font_bold)
                font_color = p_def.get("font_color")
                if font_color:
                    para_range.Font.Color.RGB = _parse_color_int(font_color)
                font_name = p_def.get("font_name")
                if font_name is not None:
                    para_range.Font.Name = str(font_name)
            paragraphs_edited += 1
        if paragraph_edits:
            result["paragraphs_edited"] = paragraphs_edited

        # Placeholder edits — target by type name
        _com_placeholder_type_map = {
            "title": 1,
            "center_title": 3,
            "subtitle": 2,
            "body": 2,
            "slide_number": 12,
            "date": 10,
            "footer": 11,
        }
        placeholders_edited = 0
        for phe in placeholder_edits:
            si = phe.get("slide_index")
            ph_type = phe.get("placeholder_type")
            if si is None or not ph_type:
                continue
            if si < 1 or si > prs.Slides.Count:
                continue
            slide = prs.Slides(si)
            target_type = _com_placeholder_type_map.get(ph_type.lower())
            if target_type is None:
                continue
            target_shape = None
            for j in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(j)
                try:
                    if shape.PlaceholderFormat is not None and shape.PlaceholderFormat.Type == target_type:
                        target_shape = shape
                        break
                except Exception:
                    continue
            if target_shape is None:
                continue
            text = phe.get("text")
            if text is not None and target_shape.HasTextFrame:
                target_shape.TextFrame.TextRange.Text = str(text)
            if target_shape.HasTextFrame:
                font_size = phe.get("font_size")
                if font_size is not None:
                    target_shape.TextFrame.TextRange.Font.Size = font_size
                font_bold = phe.get("font_bold")
                if font_bold is not None:
                    target_shape.TextFrame.TextRange.Font.Bold = bool(font_bold)
                font_color = phe.get("font_color")
                if font_color:
                    target_shape.TextFrame.TextRange.Font.Color.RGB = _parse_color_int(font_color)
                font_name = phe.get("font_name")
                if font_name is not None:
                    target_shape.TextFrame.TextRange.Font.Name = str(font_name)
            placeholders_edited += 1
        if placeholder_edits:
            result["placeholders_edited"] = placeholders_edited

        # Image replacements — swap image shapes preserving position/size
        images_replaced = 0
        for ir in image_replacements:
            si = ir.get("slide_index")
            shi = ir.get("shape_index")
            img_path = ir.get("image_path")
            if si is None or shi is None or not img_path:
                continue
            if si < 1 or si > prs.Slides.Count:
                continue
            slide = prs.Slides(si)
            if shi < 1 or shi > slide.Shapes.Count:
                continue
            shape = slide.Shapes(shi)
            # Check if shape is a picture (msoShapeType = 13)
            try:
                if shape.Type != 13:
                    continue
            except Exception:
                continue
            img_resolved, img_err = validate_path(img_path, kwargs.get("working_dir", _get_working_dir()))
            if img_err or not os.path.isfile(img_resolved):
                continue
            left = shape.Left
            top = shape.Top
            w = shape.Width
            h = shape.Height
            shape.Delete()
            slide.Shapes.AddPicture(
                FileName=os.path.abspath(img_resolved),
                LinkToFile=False,
                SaveWithDocument=True,
                Left=left,
                Top=top,
                Width=w,
                Height=h,
            )
            images_replaced += 1
        if image_replacements:
            result["images_replaced"] = images_replaced

        # Shape text/formatting edits
        shapes_edited = 0
        for se in shape_edits:
            si = se.get("slide_index")
            shi = se.get("shape_index")
            if si is None or shi is None or si < 1 or si > prs.Slides.Count:
                continue
            slide = prs.Slides(si)
            if shi < 1 or shi > slide.Shapes.Count:
                continue
            shape = slide.Shapes(shi)
            text = se.get("text")
            if text is not None and shape.HasTextFrame:
                shape.TextFrame.TextRange.Text = str(text)
            font_size = se.get("font_size")
            if font_size is not None and shape.HasTextFrame:
                shape.TextFrame.TextRange.Font.Size = font_size
            font_bold = se.get("font_bold")
            if font_bold is not None and shape.HasTextFrame:
                shape.TextFrame.TextRange.Font.Bold = bool(font_bold)
            font_color = se.get("font_color")
            if font_color is not None and shape.HasTextFrame:
                shape.TextFrame.TextRange.Font.Color.RGB = _parse_color_int(font_color)
            font_name = se.get("font_name")
            if font_name is not None and shape.HasTextFrame:
                shape.TextFrame.TextRange.Font.Name = str(font_name)
            shapes_edited += 1
        if shape_edits:
            result["shapes_edited"] = shapes_edited

        # Notes edits
        notes_edited = 0
        for ne in notes_edits:
            si = ne.get("slide_index")
            text = ne.get("text")
            if si is None or text is None or si < 1 or si > prs.Slides.Count:
                continue
            slide = prs.Slides(si)
            slide.NotesPage.Shapes(2).TextFrame.TextRange.Text = str(text)
            notes_edited += 1
        if notes_edits:
            result["notes_edited"] = notes_edited

        # Duplicate slides (before deletes, ascending order)
        slides_duplicated = 0
        for di in sorted(duplicate_slides):
            if di < 1 or di > prs.Slides.Count:
                continue
            prs.Slides(di).Duplicate()
            slides_duplicated += 1
        if duplicate_slides:
            result["slides_duplicated"] = slides_duplicated

        # Delete slides (descending order to preserve indices)
        slides_deleted = 0
        for di in sorted(delete_slides, reverse=True):
            if di < 1 or di > prs.Slides.Count:
                continue
            prs.Slides(di).Delete()
            slides_deleted += 1
        if delete_slides:
            result["slides_deleted"] = slides_deleted

        # Append new slides
        for slide_def in slides:
            _add_slide_com(prs, slide_def)
        if slides:
            result["slides_appended"] = len(slides)

        prs.Save()
    finally:
        prs.Close()

    return result


# --- transitions ---


def _transitions_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        transition = kwargs.get("transition") or {}

        if not transition:
            prs.Close()
            return {"error": "transition is required for transitions action"}

        if slide_index is not None:
            slide, serr = _get_slide_com(prs, slide_index)
            if serr:
                prs.Close()
                return {"error": serr}
            slides_to_set = [slide]
        else:
            slides_to_set = [prs.Slides(i) for i in range(1, prs.Slides.Count + 1)]

        effect = transition.get("effect")
        advance_on_time = transition.get("advance_on_time")
        advance_time = transition.get("advance_time")
        speed = transition.get("speed")

        effect_map = {
            "none": 0,
            "blinds_horizontal": 769,
            "blinds_vertical": 770,
            "checkerboard": 1025,
            "cover_down": 1284,
            "cut": 257,
            "dissolve": 1537,
            "fade": 1793,
            "push_down": 3341,
            "push_left": 3342,
            "push_right": 3343,
            "push_up": 3344,
            "random": 513,
            "split_horizontal_in": 2817,
            "wipe_down": 2052,
            "wipe_left": 2049,
            "wipe_right": 2050,
            "wipe_up": 2051,
        }

        speed_map = {"slow": 3, "medium": 2, "fast": 1}

        for slide in slides_to_set:
            sst = slide.SlideShowTransition
            if effect is not None:
                effect_val = effect_map.get(effect.lower())
                if effect_val is not None:
                    sst.EntryEffect = effect_val
            if advance_on_time is not None:
                sst.AdvanceOnTime = advance_on_time
            if advance_time is not None:
                sst.AdvanceTime = int(advance_time)
            if speed is not None:
                speed_val = speed_map.get(speed.lower())
                if speed_val is not None:
                    sst.Speed = speed_val

        prs.Save()
    finally:
        prs.Close()

    target = f"slide {slide_index}" if slide_index else "all slides"
    return {"result": f"Set transitions on {target} in {display_path}", "path": display_path}


def _transitions_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("transitions")


# --- animations ---


def _animations_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        operation = kwargs.get("operation", "add")

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        if operation == "list":
            animations: list[dict[str, Any]] = []
            try:
                seq = slide.TimeLine.MainSequence
                for i in range(1, seq.Count + 1):
                    effect = seq(i)
                    animations.append(
                        {
                            "index": i,
                            "shape_name": effect.Shape.Name if effect.Shape else None,
                            "effect_type": effect.EffectType,
                        }
                    )
            except Exception:
                pass
            prs.Close()
            return {"animations": animations, "slide_index": slide_index}

        # add operation
        shape_index: int | None = kwargs.get("shape_index")
        effect_id: int | None = kwargs.get("effect_id")

        if shape_index is None or effect_id is None:
            prs.Close()
            return {"error": "shape_index and effect_id required for animations add"}

        if shape_index < 1 or shape_index > slide.Shapes.Count:
            prs.Close()
            return {
                "error": f"shape_index {shape_index} out of range (1-{slide.Shapes.Count})",
            }

        shape = slide.Shapes(shape_index)
        slide.TimeLine.MainSequence.AddEffect(Shape=shape, effectId=effect_id)

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Added animation to shape {shape_index} on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _animations_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("animations")


# --- insert_image ---


def _insert_image_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        image_path: str | None = kwargs.get("image_path")

        if not image_path:
            prs.Close()
            return {"error": "image_path is required for insert_image"}

        img_resolved, img_err = validate_path(image_path, kwargs.get("working_dir", _get_working_dir()))
        if img_err:
            prs.Close()
            return {"error": img_err}

        if not os.path.isfile(img_resolved):
            prs.Close()
            return {"error": f"Image file not found: {image_path}"}

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        left = int(kwargs.get("left", 1) * _EMU_PER_INCH)
        top = int(kwargs.get("top", 1) * _EMU_PER_INCH)
        width = int(kwargs.get("width", 4) * _EMU_PER_INCH)
        height = int(kwargs.get("height", 3) * _EMU_PER_INCH)

        slide.Shapes.AddPicture(
            FileName=os.path.abspath(img_resolved),
            LinkToFile=False,
            SaveWithDocument=True,
            Left=left,
            Top=top,
            Width=width,
            Height=height,
        )

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Inserted image on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _insert_image_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation
    from pptx.util import Inches

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    slide_index: int | None = kwargs.get("slide_index")
    image_path: str | None = kwargs.get("image_path")

    if not image_path:
        return {"error": "image_path is required for insert_image"}

    img_resolved, img_err = validate_path(image_path, kwargs.get("working_dir", _get_working_dir()))
    if img_err:
        return {"error": img_err}

    if not os.path.isfile(img_resolved):
        return {"error": f"Image file not found: {image_path}"}

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    if slide_index is None:
        return {"error": "slide_index is required for insert_image"}

    slides_list = list(prs.slides)
    if slide_index < 1 or slide_index > len(slides_list):
        return {"error": f"slide_index {slide_index} out of range (1-{len(slides_list)})"}

    slide = slides_list[slide_index - 1]

    left = Inches(kwargs.get("left", 1))
    top = Inches(kwargs.get("top", 1))
    width = Inches(kwargs.get("width", 4))
    height = Inches(kwargs.get("height", 3))

    slide.shapes.add_picture(img_resolved, left, top, width, height)

    prs.save(resolved)
    return {
        "result": f"Inserted image on slide {slide_index} in {display_path}",
        "path": display_path,
    }


# --- insert_shape ---


def _insert_shape_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        shape_type: str | None = kwargs.get("shape_type")

        if not shape_type:
            prs.Close()
            return {"error": "shape_type is required for insert_shape"}

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        mso_type = _COM_SHAPE_TYPE_MAP.get(shape_type.lower())
        if mso_type is None:
            prs.Close()
            return {
                "error": f"Unknown shape_type: {shape_type}. Available: {', '.join(_SHAPE_TYPE_NAMES)}",
            }

        left = int(kwargs.get("left", 2) * _EMU_PER_INCH)
        top = int(kwargs.get("top", 2) * _EMU_PER_INCH)
        width = int(kwargs.get("width", 3) * _EMU_PER_INCH)
        height = int(kwargs.get("height", 2) * _EMU_PER_INCH)

        slide.Shapes.AddShape(Type=mso_type, Left=left, Top=top, Width=width, Height=height)

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Inserted {shape_type} on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _insert_shape_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    slide_index: int | None = kwargs.get("slide_index")
    shape_type: str | None = kwargs.get("shape_type")

    if not shape_type:
        return {"error": "shape_type is required for insert_shape"}

    # Map uses shared _SHAPE_TYPE_NAMES for consistency with COM backend
    _lib_shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "left_arrow": MSO_SHAPE.LEFT_ARROW,
        "diamond": MSO_SHAPE.DIAMOND,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "star": MSO_SHAPE.STAR_5_POINT,
    }

    normalized = shape_type.lower()
    if normalized not in _SHAPE_TYPE_NAMES:
        return {
            "error": f"Unknown shape_type: {shape_type}. Available: {', '.join(_SHAPE_TYPE_NAMES)}",
        }

    mso_shape = _lib_shape_map.get(normalized)
    if mso_shape is None:
        return {
            "error": f"Unknown shape_type: {shape_type}. Available: {', '.join(_SHAPE_TYPE_NAMES)}",
        }

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    if slide_index is None:
        return {"error": "slide_index is required for insert_shape"}

    slides_list = list(prs.slides)
    if slide_index < 1 or slide_index > len(slides_list):
        return {"error": f"slide_index {slide_index} out of range (1-{len(slides_list)})"}

    slide = slides_list[slide_index - 1]

    left = Inches(kwargs.get("left", 2))
    top = Inches(kwargs.get("top", 2))
    width = Inches(kwargs.get("width", 3))
    height = Inches(kwargs.get("height", 2))

    slide.shapes.add_shape(mso_shape, left, top, width, height)

    prs.save(resolved)
    return {
        "result": f"Inserted {shape_type} on slide {slide_index} in {display_path}",
        "path": display_path,
    }


# --- format_shape ---


def _format_shape_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        shape_index: int | None = kwargs.get("shape_index")
        fmt = kwargs.get("format") or {}

        if not fmt:
            prs.Close()
            return {"error": "format is required for format_shape"}

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        if shape_index is None:
            prs.Close()
            return {"error": "shape_index is required for format_shape"}

        if shape_index < 1 or shape_index > slide.Shapes.Count:
            prs.Close()
            return {
                "error": f"shape_index {shape_index} out of range (1-{slide.Shapes.Count})",
            }

        shape = slide.Shapes(shape_index)

        fill_color = fmt.get("fill_color")
        if fill_color:
            shape.Fill.Visible = True
            shape.Fill.ForeColor.RGB = _parse_color_int(fill_color)

        line_color = fmt.get("line_color")
        if line_color:
            shape.Line.Visible = True
            shape.Line.ForeColor.RGB = _parse_color_int(line_color)

        line_width = fmt.get("line_width")
        if line_width is not None:
            shape.Line.Weight = line_width

        shadow = fmt.get("shadow")
        if shadow is not None:
            shape.Shadow.Visible = shadow

        text = fmt.get("text")
        if text is not None and shape.HasTextFrame:
            shape.TextFrame.TextRange.Text = str(text)

        font_size = fmt.get("font_size")
        if font_size is not None and shape.HasTextFrame:
            shape.TextFrame.TextRange.Font.Size = font_size

        font_bold = fmt.get("font_bold")
        if font_bold is not None and shape.HasTextFrame:
            shape.TextFrame.TextRange.Font.Bold = font_bold

        font_color = fmt.get("font_color")
        if font_color and shape.HasTextFrame:
            shape.TextFrame.TextRange.Font.Color.RGB = _parse_color_int(font_color)

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Formatted shape {shape_index} on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _format_shape_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation
    from pptx.util import Pt

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    slide_index: int | None = kwargs.get("slide_index")
    shape_index: int | None = kwargs.get("shape_index")
    fmt = kwargs.get("format") or {}

    if not fmt:
        return {"error": "format is required for format_shape"}

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    if slide_index is None:
        return {"error": "slide_index is required for format_shape"}

    slides_list = list(prs.slides)
    if slide_index < 1 or slide_index > len(slides_list):
        return {"error": f"slide_index {slide_index} out of range (1-{len(slides_list)})"}

    slide = slides_list[slide_index - 1]

    if shape_index is None:
        return {"error": "shape_index is required for format_shape"}

    shapes_list = list(slide.shapes)
    if shape_index < 1 or shape_index > len(shapes_list):
        return {"error": f"shape_index {shape_index} out of range (1-{len(shapes_list)})"}

    shape = shapes_list[shape_index - 1]

    fill_color = fmt.get("fill_color")
    if fill_color:
        parsed = _parse_rgb_color(fill_color)
        if parsed is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = parsed

    line_color = fmt.get("line_color")
    if line_color:
        parsed = _parse_rgb_color(line_color)
        if parsed is not None:
            shape.line.color.rgb = parsed

    line_width = fmt.get("line_width")
    if line_width is not None:
        shape.line.width = Pt(line_width)

    text = fmt.get("text")
    if text is not None and shape.has_text_frame:
        shape.text_frame.text = str(text)

    font_size = fmt.get("font_size")
    font_bold = fmt.get("font_bold")
    font_color = fmt.get("font_color")

    if (font_size is not None or font_bold is not None or font_color) and shape.has_text_frame:
        parsed_font_color = _parse_rgb_color(font_color) if font_color else None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if font_size is not None:
                    run.font.size = Pt(font_size)
                if font_bold is not None:
                    run.font.bold = font_bold
                if parsed_font_color is not None:
                    run.font.color.rgb = parsed_font_color

    prs.save(resolved)
    return {
        "result": f"Formatted shape {shape_index} on slide {slide_index} in {display_path}",
        "path": display_path,
    }


# --- master_layout ---


def _master_layout_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path, read_only=True)
    if err:
        return {"error": err}

    try:
        operation = kwargs.get("operation", "list")

        if operation == "list":
            layouts: list[dict[str, Any]] = []
            master = prs.SlideMaster
            for i in range(1, master.CustomLayouts.Count + 1):
                layout = master.CustomLayouts(i)
                layouts.append({"index": i, "name": layout.Name})
            prs.Close()
            return {"layouts": layouts, "count": len(layouts)}

        # apply operation requires a writable presentation
        prs.Close()

        if operation != "apply":
            return {"error": f"Unknown operation: {operation}. Use 'list' or 'apply'"}

        _, prs2, err2 = _open_pres_com(manager, resolved, display_path)
        if err2:
            return {"error": err2}

        try:
            slide_index: int | None = kwargs.get("slide_index")
            layout_index: int | None = kwargs.get("layout_index")

            if slide_index is None or layout_index is None:
                prs2.Close()
                return {"error": "slide_index and layout_index required for apply"}

            slide, serr = _get_slide_com(prs2, slide_index)
            if serr:
                prs2.Close()
                return {"error": serr}

            master = prs2.SlideMaster
            if layout_index < 1 or layout_index > master.CustomLayouts.Count:
                prs2.Close()
                return {
                    "error": f"layout_index {layout_index} out of range (1-{master.CustomLayouts.Count})",
                }

            new_layout = master.CustomLayouts(layout_index)
            slide.CustomLayout = new_layout
            prs2.Save()
        finally:
            prs2.Close()

        return {
            "result": f"Applied layout {layout_index} to slide {slide_index} in {display_path}",
            "path": display_path,
        }

    except Exception:
        try:
            prs.Close()
        except Exception:
            pass
        raise


def _master_layout_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    operation = kwargs.get("operation", "list")

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    if operation == "list":
        layouts: list[dict[str, Any]] = []
        for master in prs.slide_masters:
            for i, layout in enumerate(master.slide_layouts, 1):
                layouts.append({"index": i, "name": layout.name})
        return {"layouts": layouts, "count": len(layouts)}

    if operation == "apply":
        return _com_only_error("master_layout apply")

    return {"error": f"Unknown operation: {operation}. Use 'list' or 'apply'"}


# --- reorder_slides ---


def _reorder_slides_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        operation = kwargs.get("operation", "move")
        slide_index: int | None = kwargs.get("slide_index")

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        if operation == "delete":
            slide.Delete()
            prs.Save()
            prs.Close()
            return {
                "result": f"Deleted slide {slide_index} from {display_path}",
                "path": display_path,
            }

        if operation == "duplicate":
            slide.Duplicate()
            prs.Save()
            prs.Close()
            return {
                "result": f"Duplicated slide {slide_index} in {display_path}",
                "path": display_path,
            }

        if operation == "move":
            new_position: int | None = kwargs.get("new_position")
            if new_position is None:
                prs.Close()
                return {"error": "new_position is required for move operation"}
            if new_position < 1 or new_position > prs.Slides.Count:
                prs.Close()
                return {
                    "error": f"new_position {new_position} out of range (1-{prs.Slides.Count})",
                }
            slide.MoveTo(new_position)
            prs.Save()
            prs.Close()
            return {
                "result": f"Moved slide {slide_index} to position {new_position} in {display_path}",
                "path": display_path,
            }

        prs.Close()
        return {"error": f"Unknown operation: {operation}. Use 'move', 'delete', or 'duplicate'"}

    except Exception:
        try:
            prs.Close()
        except Exception:
            pass
        raise


def _reorder_slides_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    operation = kwargs.get("operation", "move")
    slide_index: int | None = kwargs.get("slide_index")

    if slide_index is None:
        return {"error": "slide_index is required for reorder_slides"}

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    slides_list = list(prs.slides)
    if slide_index < 1 or slide_index > len(slides_list):
        return {"error": f"slide_index {slide_index} out of range (1-{len(slides_list)})"}

    if operation == "delete":
        # Delete slide via XML manipulation.
        # Also drop the relationship to avoid orphaning the slide XML part.
        rns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        id_list = prs.element.find(".//p:sldIdLst", ns)

        if id_list is not None:
            children = list(id_list)
            if slide_index - 1 < len(children):
                sld_id_elem = children[slide_index - 1]
                r_id = sld_id_elem.get(f"{{{rns}}}id")
                id_list.remove(sld_id_elem)
                # Drop the relationship so the slide part is not orphaned
                if r_id is not None:
                    try:
                        prs.part.drop_rel(r_id)
                    except (AttributeError, KeyError):
                        pass  # Relationship cleanup not supported in this python-pptx version

        prs.save(resolved)
        return {
            "result": f"Deleted slide {slide_index} from {display_path}",
            "path": display_path,
        }

    if operation == "duplicate":
        return _com_only_error("reorder_slides duplicate")

    if operation == "move":
        new_position: int | None = kwargs.get("new_position")
        if new_position is None:
            return {"error": "new_position is required for move operation"}
        if new_position < 1 or new_position > len(slides_list):
            return {
                "error": f"new_position {new_position} out of range (1-{len(slides_list)})",
            }

        ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        id_list = prs.element.find(".//p:sldIdLst", ns)

        if id_list is not None:
            children = list(id_list)
            if slide_index - 1 < len(children):
                elem = children[slide_index - 1]
                id_list.remove(elem)
                id_list.insert(new_position - 1, elem)

        prs.save(resolved)
        return {
            "result": f"Moved slide {slide_index} to position {new_position} in {display_path}",
            "path": display_path,
        }

    return {"error": f"Unknown operation: {operation}. Use 'move', 'delete', or 'duplicate'"}


# --- embed_chart ---


def _embed_chart_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        chart_type: str | None = kwargs.get("chart_type", "column")
        chart_title: str | None = kwargs.get("chart_title")
        data: list[dict[str, Any]] | None = kwargs.get("data")

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        # COM xlChartType constants
        type_map = {
            "bar": 57,
            "line": 4,
            "pie": 5,
            "column": 51,
            "area": 1,
            "scatter": -4169,
        }

        xl_type = type_map.get((chart_type or "column").lower(), 51)

        left = int(kwargs.get("left", 1) * _EMU_PER_INCH)
        top = int(kwargs.get("top", 1) * _EMU_PER_INCH)
        width = int(kwargs.get("width", 8) * _EMU_PER_INCH)
        height = int(kwargs.get("height", 5) * _EMU_PER_INCH)

        chart_shape = slide.Shapes.AddChart2(
            Style=-1,
            Type=xl_type,
            Left=left,
            Top=top,
            Width=width,
            Height=height,
        )
        chart = chart_shape.Chart

        if chart_title:
            chart.HasTitle = True
            chart.ChartTitle.Text = str(chart_title)

        if data:
            # Populate chart data via the chart's data sheet
            try:
                ws = chart.ChartData.Workbook.Worksheets(1)
                for col_idx, series in enumerate(data):
                    name = series.get("name", f"Series {col_idx + 1}")
                    values = series.get("values", [])
                    ws.Cells(1, col_idx + 2).Value = name
                    for row_idx, val in enumerate(values):
                        ws.Cells(row_idx + 2, col_idx + 2).Value = val
                        if col_idx == 0:
                            ws.Cells(row_idx + 2, 1).Value = f"Cat {row_idx + 1}"
            except Exception:
                pass

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Inserted chart on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _embed_chart_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("embed_chart")


# --- embed_table ---


def _embed_table_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        rows: int | None = kwargs.get("rows")
        cols: int | None = kwargs.get("cols")
        data: list[list[Any]] | None = kwargs.get("data")

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        if data:
            rows = rows or len(data)
            cols = cols or (len(data[0]) if data else 1)

        if not rows or not cols:
            prs.Close()
            return {"error": "rows and cols (or data) required for embed_table"}

        if rows < 1 or cols < 1:
            prs.Close()
            return {"error": "rows and cols must be positive integers"}

        tbl_width = kwargs.get("width", 8)
        tbl_height = kwargs.get("height", 3)
        if tbl_width <= 0 or tbl_height <= 0:
            prs.Close()
            return {"error": "width and height must be positive numbers"}

        left = int(kwargs.get("left", 1) * _EMU_PER_INCH)
        top = int(kwargs.get("top", 2) * _EMU_PER_INCH)
        width = int(tbl_width * _EMU_PER_INCH)
        height = int(tbl_height * _EMU_PER_INCH)

        table_shape = slide.Shapes.AddTable(
            NumRows=rows,
            NumColumns=cols,
            Left=left,
            Top=top,
            Width=width,
            Height=height,
        )
        table = table_shape.Table

        if data:
            for r_idx, row_data in enumerate(data):
                if r_idx >= rows:
                    break
                for c_idx, val in enumerate(row_data):
                    if c_idx >= cols:
                        break
                    table.Cell(r_idx + 1, c_idx + 1).Shape.TextFrame.TextRange.Text = str(val)

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Inserted {rows}x{cols} table on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _embed_table_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation
    from pptx.util import Inches

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    slide_index: int | None = kwargs.get("slide_index")
    rows: int | None = kwargs.get("rows")
    cols: int | None = kwargs.get("cols")
    data: list[list[Any]] | None = kwargs.get("data")

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    if slide_index is None:
        return {"error": "slide_index is required for embed_table"}

    slides_list = list(prs.slides)
    if slide_index < 1 or slide_index > len(slides_list):
        return {"error": f"slide_index {slide_index} out of range (1-{len(slides_list)})"}

    slide = slides_list[slide_index - 1]

    if data:
        rows = rows or len(data)
        cols = cols or (len(data[0]) if data else 1)

    if not rows or not cols:
        return {"error": "rows and cols (or data) required for embed_table"}

    if rows < 1 or cols < 1:
        return {"error": "rows and cols must be positive integers"}

    tbl_width = kwargs.get("width", 8)
    tbl_height = kwargs.get("height", 3)
    if tbl_width <= 0 or tbl_height <= 0:
        return {"error": "width and height must be positive numbers"}

    left = Inches(kwargs.get("left", 1))
    top = Inches(kwargs.get("top", 2))
    width = Inches(tbl_width)
    height = Inches(tbl_height)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if data:
        for r_idx, row_data in enumerate(data):
            if r_idx >= rows:
                break
            for c_idx, val in enumerate(row_data):
                if c_idx >= cols:
                    break
                table.cell(r_idx, c_idx).text = str(val)

    prs.save(resolved)
    return {
        "result": f"Inserted {rows}x{cols} table on slide {slide_index} in {display_path}",
        "path": display_path,
    }


# --- export_pdf ---


def _export_pdf_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        output_path = kwargs.get("output_path")
        if not output_path:
            output_path = os.path.splitext(resolved)[0] + ".pdf"
        else:
            out_resolved, out_err = validate_path(output_path, kwargs.get("working_dir", _get_working_dir()))
            if out_err:
                prs.Close()
                return {"error": out_err}
            output_path = out_resolved

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        # ppFixedFormatTypePDF = 32
        prs.ExportAsFixedFormat(Path=os.path.abspath(output_path), FixedFormatType=32)
    finally:
        prs.Close()

    return {"result": f"Exported PDF to {output_path}", "path": display_path}


def _export_pdf_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("export_pdf")


# --- hyperlinks ---


def _hyperlinks_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        shape_index: int | None = kwargs.get("shape_index")
        url: str | None = kwargs.get("url")
        display_text_val: str | None = kwargs.get("display_text")
        operation = kwargs.get("operation", "add")

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        if operation == "list":
            links: list[dict[str, Any]] = []
            for i in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(i)
                if shape.HasTextFrame:
                    tr = shape.TextFrame.TextRange
                    try:
                        hl = tr.ActionSettings(1).Hyperlink
                        if hl and hl.Address:
                            links.append(
                                {
                                    "shape_index": i,
                                    "shape_name": shape.Name,
                                    "url": hl.Address,
                                    "text": tr.Text,
                                }
                            )
                    except Exception:
                        pass
            prs.Close()
            return {"hyperlinks": links, "slide_index": slide_index}

        # add operation
        if not url:
            prs.Close()
            return {"error": "url is required for hyperlinks add"}

        allowed_url_schemes = ("http://", "https://", "mailto:")
        if not url.lower().startswith(allowed_url_schemes):
            prs.Close()
            return {"error": "url must start with http://, https://, or mailto:"}

        if shape_index is None:
            prs.Close()
            return {"error": "shape_index is required for hyperlinks add"}

        if shape_index < 1 or shape_index > slide.Shapes.Count:
            prs.Close()
            return {
                "error": f"shape_index {shape_index} out of range (1-{slide.Shapes.Count})",
            }

        shape = slide.Shapes(shape_index)
        if shape.HasTextFrame:
            tr = shape.TextFrame.TextRange
            # ppMouseClick = 1
            tr.ActionSettings(1).Hyperlink.Address = url
            if display_text_val:
                tr.Text = display_text_val

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Added hyperlink to shape {shape_index} on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _hyperlinks_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    slide_index: int | None = kwargs.get("slide_index")
    shape_index: int | None = kwargs.get("shape_index")
    url: str | None = kwargs.get("url")
    display_text_val: str | None = kwargs.get("display_text")
    operation = kwargs.get("operation", "add")

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    if slide_index is None:
        return {"error": "slide_index is required for hyperlinks"}

    slides_list = list(prs.slides)
    if slide_index < 1 or slide_index > len(slides_list):
        return {"error": f"slide_index {slide_index} out of range (1-{len(slides_list)})"}

    slide = slides_list[slide_index - 1]

    if operation == "list":
        links: list[dict[str, Any]] = []
        for i, shape in enumerate(slide.shapes, 1):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            links.append(
                                {
                                    "shape_index": i,
                                    "shape_name": shape.name,
                                    "url": run.hyperlink.address,
                                    "text": run.text,
                                }
                            )
        return {"hyperlinks": links, "slide_index": slide_index}

    # add operation
    if not url:
        return {"error": "url is required for hyperlinks add"}

    allowed_url_schemes = ("http://", "https://", "mailto:")
    if not url.lower().startswith(allowed_url_schemes):
        return {"error": "url must start with http://, https://, or mailto:"}

    if shape_index is None:
        return {"error": "shape_index is required for hyperlinks add"}

    shapes_list = list(slide.shapes)
    if shape_index < 1 or shape_index > len(shapes_list):
        return {"error": f"shape_index {shape_index} out of range (1-{len(shapes_list)})"}

    shape = shapes_list[shape_index - 1]

    if shape.has_text_frame:
        if display_text_val:
            shape.text_frame.text = display_text_val
        applied = False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.hyperlink.address = url
                applied = True
                break
            if applied:
                break
        if not applied:
            return {"error": "Shape has no text runs to apply hyperlink to"}
    else:
        return {"error": "Shape has no text frame to apply hyperlink to"}

    prs.save(resolved)
    return {
        "result": f"Added hyperlink to shape {shape_index} on slide {slide_index} in {display_path}",
        "path": display_path,
    }


# --- headers_footers ---


def _headers_footers_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        footer_text: str | None = kwargs.get("footer_text")
        slide_numbers: bool | None = kwargs.get("slide_numbers")
        date_time: str | None = kwargs.get("date_time")
        operation = kwargs.get("operation", "set")

        if operation == "list":
            info: dict[str, Any] = {}
            try:
                master_hf = prs.SlideMaster.HeadersFooters
                info["master_footer"] = master_hf.Footer.Text if master_hf.Footer.Visible else None
                info["master_slide_number"] = master_hf.SlideNumber.Visible
                info["master_date_time"] = master_hf.DateAndTime.Visible
            except Exception:
                info["master_footer"] = None
            prs.Close()
            return {"headers_footers": info}

        # set operation
        if slide_index is not None:
            slide, serr = _get_slide_com(prs, slide_index)
            if serr:
                prs.Close()
                return {"error": serr}
            hf = slide.HeadersFooters
        else:
            hf = prs.SlideMaster.HeadersFooters

        if footer_text is not None:
            hf.Footer.Visible = True
            hf.Footer.Text = str(footer_text)

        if slide_numbers is not None:
            hf.SlideNumber.Visible = slide_numbers

        if date_time is not None:
            hf.DateAndTime.Visible = True
            hf.DateAndTime.UseFormat = False
            hf.DateAndTime.Text = str(date_time)

        prs.Save()
    finally:
        prs.Close()

    target = f"slide {slide_index}" if slide_index else "presentation master"
    return {
        "result": f"Updated headers/footers on {target} in {display_path}",
        "path": display_path,
    }


def _headers_footers_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    slide_index: int | None = kwargs.get("slide_index")
    footer_text: str | None = kwargs.get("footer_text")
    slide_numbers: bool | None = kwargs.get("slide_numbers")
    date_time: str | None = kwargs.get("date_time")
    operation = kwargs.get("operation", "set")

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    if operation == "list":
        # Read placeholder content from slides
        info: dict[str, Any] = {"slides": []}
        for i, slide in enumerate(prs.slides, 1):
            slide_info: dict[str, Any] = {"slide_index": i}
            for ph in slide.placeholders:
                # 12 = slide number, 13 = date, 14 = footer
                if ph.placeholder_format.idx == 14:
                    slide_info["footer"] = ph.text
                elif ph.placeholder_format.idx == 12:
                    slide_info["slide_number"] = True
                elif ph.placeholder_format.idx == 13:
                    slide_info["date_time"] = ph.text
            info["slides"].append(slide_info)
        return {"headers_footers": info}

    # set operation — manipulate XML for footer/slide number/date placeholders
    if slide_index is not None:
        slides_list = list(prs.slides)
        if slide_index < 1 or slide_index > len(slides_list):
            return {"error": f"slide_index {slide_index} out of range (1-{len(slides_list)})"}
        target_slides = [slides_list[slide_index - 1]]
    else:
        target_slides = list(prs.slides)

    for slide in target_slides:
        for ph in slide.placeholders:
            if footer_text is not None and ph.placeholder_format.idx == 14:
                ph.text = str(footer_text)
            if date_time is not None and ph.placeholder_format.idx == 13:
                ph.text = str(date_time)

        # For slide numbers, we can toggle visibility via the sp element
        if slide_numbers is not None:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 12:
                    sp = ph._element
                    if not slide_numbers:
                        sp.getparent().remove(sp)

    prs.save(resolved)
    target = f"slide {slide_index}" if slide_index else "all slides"
    return {
        "result": f"Updated headers/footers on {target} in {display_path}",
        "path": display_path,
    }


# --- sections ---


def _sections_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        operation = kwargs.get("operation", "list")

        if operation == "list":
            sections: list[dict[str, Any]] = []
            try:
                sp = prs.SectionProperties
                for i in range(1, sp.Count + 1):
                    sections.append(
                        {
                            "index": i,
                            "name": sp.Name(i),
                            "first_slide": sp.FirstSlide(i),
                            "slide_count": sp.SlidesCount(i),
                        }
                    )
            except Exception:
                pass
            prs.Close()
            return {"sections": sections, "count": len(sections)}

        if operation == "add":
            section_name: str | None = kwargs.get("section_name")
            slide_index: int | None = kwargs.get("slide_index")

            if not section_name:
                prs.Close()
                return {"error": "section_name is required for sections add"}

            idx = slide_index if slide_index is not None else prs.Slides.Count
            prs.SectionProperties.AddSection(idx, section_name)
            prs.Save()
            prs.Close()
            return {
                "result": f"Added section '{section_name}' at slide {idx} in {display_path}",
                "path": display_path,
            }

        if operation == "delete":
            # Accept section_index or fall back to slide_index for backward compat
            section_index: int | None = kwargs.get("section_index") or kwargs.get("slide_index")
            if section_index is None:
                prs.Close()
                return {"error": "section_index is required for sections delete"}
            # DeleteSection(index, deleteSlides)
            prs.SectionProperties.Delete(section_index, False)
            prs.Save()
            prs.Close()
            return {
                "result": f"Deleted section {section_index} from {display_path}",
                "path": display_path,
            }

        prs.Close()
        return {"error": f"Unknown operation: {operation}. Use 'list', 'add', or 'delete'"}

    except Exception:
        try:
            prs.Close()
        except Exception:
            pass
        raise


def _sections_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("sections")


# --- group_shapes ---


def _group_shapes_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        shape_indices: list[int] | None = kwargs.get("shape_indices")
        operation = kwargs.get("operation", "add")

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        if operation == "add":
            if not shape_indices or len(shape_indices) < 2:
                prs.Close()
                return {"error": "shape_indices (at least 2) required for group_shapes"}

            for idx in shape_indices:
                if idx < 1 or idx > slide.Shapes.Count:
                    prs.Close()
                    return {
                        "error": f"shape_index {idx} out of range (1-{slide.Shapes.Count})",
                    }

            # Build an array of shape names for Range()
            shape_names = [slide.Shapes(i).Name for i in shape_indices]
            shape_range = slide.Shapes.Range(shape_names)
            shape_range.Group()

            prs.Save()
            prs.Close()
            return {
                "result": f"Grouped shapes {shape_indices} on slide {slide_index} in {display_path}",
                "path": display_path,
            }

        if operation == "delete":
            # Ungroup
            shape_index_val: int | None = kwargs.get("shape_index")
            if shape_index_val is None:
                prs.Close()
                return {"error": "shape_index required for ungroup (delete) operation"}

            if shape_index_val < 1 or shape_index_val > slide.Shapes.Count:
                prs.Close()
                return {
                    "error": f"shape_index {shape_index_val} out of range (1-{slide.Shapes.Count})",
                }

            slide.Shapes(shape_index_val).Ungroup()
            prs.Save()
            prs.Close()
            return {
                "result": f"Ungrouped shape {shape_index_val} on slide {slide_index} in {display_path}",
                "path": display_path,
            }

        prs.Close()
        return {"error": f"Unknown operation: {operation}. Use 'add' (group) or 'delete' (ungroup)"}

    except Exception:
        try:
            prs.Close()
        except Exception:
            pass
        raise


def _group_shapes_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("group_shapes")


# --- audio_video ---


def _audio_video_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        media_path: str | None = kwargs.get("media_path")

        if not media_path:
            prs.Close()
            return {"error": "media_path is required for audio_video"}

        media_resolved, media_err = validate_path(media_path, kwargs.get("working_dir", _get_working_dir()))
        if media_err:
            prs.Close()
            return {"error": media_err}

        if not os.path.isfile(media_resolved):
            prs.Close()
            return {"error": f"Media file not found: {media_path}"}

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        left = int(kwargs.get("left", 2) * _EMU_PER_INCH)
        top = int(kwargs.get("top", 2) * _EMU_PER_INCH)
        width = int(kwargs.get("width", 3) * _EMU_PER_INCH)
        height = int(kwargs.get("height", 2) * _EMU_PER_INCH)

        slide.Shapes.AddMediaObject2(
            FileName=os.path.abspath(media_resolved),
            Left=left,
            Top=top,
            Width=width,
            Height=height,
        )

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Inserted media on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _audio_video_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("audio_video")


# --- smartart ---


def _smartart_com(manager: Any, resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    _, prs, err = _open_pres_com(manager, resolved, display_path)
    if err:
        return {"error": err}

    try:
        slide_index: int | None = kwargs.get("slide_index")
        smartart_layout: str | None = kwargs.get("smartart_layout")
        smartart_items: list[str] | None = kwargs.get("smartart_items")

        slide, serr = _get_slide_com(prs, slide_index)
        if serr:
            prs.Close()
            return {"error": serr}

        left = int(kwargs.get("left", 1) * _EMU_PER_INCH)
        top = int(kwargs.get("top", 1) * _EMU_PER_INCH)
        width = int(kwargs.get("width", 8) * _EMU_PER_INCH)
        height = int(kwargs.get("height", 5) * _EMU_PER_INCH)

        # SmartArt requires the SmartArt layout object
        # Find layout by name from Application.SmartArtLayouts
        ppt = prs.Application
        sa_layout = None
        layout_name = (smartart_layout or "basic_block_list").lower().replace("_", " ")

        try:
            for i in range(1, ppt.SmartArtLayouts.Count + 1):
                if layout_name in ppt.SmartArtLayouts(i).Name.lower():
                    sa_layout = ppt.SmartArtLayouts(i)
                    break
        except Exception:
            pass

        if sa_layout is None:
            # Fall back to first available layout
            try:
                sa_layout = ppt.SmartArtLayouts(1)
            except Exception:
                prs.Close()
                return {"error": "No SmartArt layouts available"}

        smart_shape = slide.Shapes.AddSmartArt(
            Layout=sa_layout,
            Left=left,
            Top=top,
            Width=width,
            Height=height,
        )

        if smartart_items:
            try:
                sa = smart_shape.SmartArt
                nodes = sa.AllNodes
                for i, item_text in enumerate(smartart_items):
                    if i < nodes.Count:
                        nodes(i + 1).TextFrame2.TextRange.Text = str(item_text)
                    else:
                        sa.AllNodes.Add().TextFrame2.TextRange.Text = str(item_text)
            except Exception:
                pass

        prs.Save()
    finally:
        prs.Close()

    return {
        "result": f"Inserted SmartArt on slide {slide_index} in {display_path}",
        "path": display_path,
    }


def _smartart_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    return _com_only_error("smartart")


# ---------------------------------------------------------------------------
# Library backend (python-pptx)
# ---------------------------------------------------------------------------


def _add_slide_lib(prs: Any, slide_def: dict[str, Any]) -> None:
    """Add a single slide to a python-pptx presentation."""
    layout_idx = slide_def.get("layout", 1)
    try:
        layout = prs.slide_layouts[layout_idx]
    except IndexError:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    title_text = slide_def.get("title")
    if title_text and slide.shapes.title:
        slide.shapes.title.text = str(title_text)

    content = slide_def.get("content")
    bullets = slide_def.get("bullets")

    if content or bullets:
        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_placeholder = shape
                break

        if body_placeholder is not None and hasattr(body_placeholder, "text_frame"):
            tf = body_placeholder.text_frame
            if content:
                tf.text = str(content)
            elif bullets:
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)

    notes_text = slide_def.get("notes")
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = str(notes_text)


def _create_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation

    slides: list[dict[str, Any]] = kwargs.get("slides") or []
    if not slides:
        return {"error": "slides is required for create action"}
    if len(slides) > _MAX_SLIDES:
        return {"error": f"Too many slides (max {_MAX_SLIDES})"}

    prs = _Presentation()

    for slide_def in slides:
        _add_slide_lib(prs, slide_def)

    os.makedirs(os.path.dirname(resolved) or ".", exist_ok=True)
    prs.save(resolved)
    return {"result": f"Created {display_path}", "path": display_path, "slides_created": len(slides)}


def _read_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    output_parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        shape_list = list(slide.shapes)
        output_parts.append(f"--- Slide {i} ({len(shape_list)} shapes) ---")
        table_idx = 0
        for j, shape in enumerate(shape_list, 1):
            if shape.has_table:
                table_idx += 1
                tbl = shape.table
                output_parts.append(f"[Table {table_idx}: {len(tbl.rows)} rows x {len(tbl.columns)} cols]")
                for r, row in enumerate(tbl.rows):
                    if r >= 50:
                        output_parts.append("  ... (rows truncated)")
                        break
                    row_vals = [tbl.cell(r, c).text.strip() for c in range(len(tbl.columns))]
                    output_parts.append(f"  Row {r}: {row_vals}")
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    output_parts.append(f"[Shape {j}] {text}")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                output_parts.append(f"[Notes] {notes}")

    content = "\n".join(output_parts)
    if len(content) > _MAX_OUTPUT:
        content = content[:_MAX_OUTPUT] + "\n... (truncated)"

    return {
        "content": content,
        "slides": len(prs.slides),
    }


def _edit_lib(resolved: str, display_path: str, **kwargs: Any) -> dict[str, Any]:
    from pptx import Presentation as _Presentation
    from pptx.util import Pt

    if not os.path.isfile(resolved):
        return {"error": f"File not found: {display_path}"}

    try:
        prs = _Presentation(resolved)
    except Exception as exc:
        return {"error": f"Unable to read PPTX file: {display_path} ({type(exc).__name__}: {exc})"}

    replacements: list[dict[str, str]] = kwargs.get("replacements") or []
    slides_to_add: list[dict[str, Any]] = kwargs.get("slides") or []
    table_edits: list[dict[str, Any]] = kwargs.get("table_edits") or []
    shape_edits: list[dict[str, Any]] = kwargs.get("shape_edits") or []
    notes_edits: list[dict[str, Any]] = kwargs.get("notes_edits") or []
    delete_slides: list[int] = kwargs.get("delete_slides") or []
    duplicate_slides: list[int] = kwargs.get("duplicate_slides") or []
    template_fill: dict[str, str] = kwargs.get("template_fill") or {}
    table_format: list[dict[str, Any]] = kwargs.get("table_format") or []
    paragraph_edits: list[dict[str, Any]] = kwargs.get("paragraph_edits") or []
    placeholder_edits: list[dict[str, Any]] = kwargs.get("placeholder_edits") or []
    image_replacements: list[dict[str, Any]] = kwargs.get("image_replacements") or []

    has_work = (
        replacements
        or slides_to_add
        or table_edits
        or shape_edits
        or notes_edits
        or delete_slides
        or duplicate_slides
        or template_fill
        or table_format
        or paragraph_edits
        or placeholder_edits
        or image_replacements
    )
    if not has_work:
        return {
            "error": (
                "Provide at least one of: 'replacements', 'slides', 'table_edits', "
                "'shape_edits', 'notes_edits', 'delete_slides', 'duplicate_slides', "
                "'template_fill', 'table_format', 'paragraph_edits', "
                "'placeholder_edits', 'image_replacements'"
            )
        }

    if slides_to_add and len(slides_to_add) > _MAX_SLIDES:
        return {"error": f"Too many slides to append (max {_MAX_SLIDES})"}

    for _arr_name, _arr in [
        ("table_edits", table_edits),
        ("table_format", table_format),
        ("shape_edits", shape_edits),
        ("notes_edits", notes_edits),
        ("paragraph_edits", paragraph_edits),
        ("placeholder_edits", placeholder_edits),
        ("image_replacements", image_replacements),
        ("delete_slides", delete_slides),
        ("duplicate_slides", duplicate_slides),
    ]:
        if len(_arr) > _MAX_EDIT_OPS:
            return {"error": f"Too many {_arr_name} entries (max {_MAX_EDIT_OPS})"}

    if len(template_fill) > _MAX_EDIT_OPS:
        return {"error": f"Too many template_fill keys (max {_MAX_EDIT_OPS})"}

    current_count = len(prs.slides)
    if slides_to_add and current_count + len(slides_to_add) > _MAX_SLIDES:
        return {"error": f"Total slides would exceed limit (max {_MAX_SLIDES})"}

    slide_list = list(prs.slides)
    result: dict[str, Any] = {"result": f"Edited {display_path}", "path": display_path}

    # Template fill — {{key}} replacement throughout entire presentation
    tokens_replaced = 0
    if template_fill:
        for key, value in template_fill.items():
            token = "{{" + str(key) + "}}"
            val_str = str(value)
            for slide in slide_list:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if token in run.text:
                                    run.text = run.text.replace(token, val_str)
                                    tokens_replaced += 1
                    if shape.has_table:
                        tbl = shape.table
                        for r in range(len(tbl.rows)):
                            for c in range(len(tbl.columns)):
                                cell = tbl.cell(r, c)
                                if cell.text_frame:
                                    for para in cell.text_frame.paragraphs:
                                        for run in para.runs:
                                            if token in run.text:
                                                run.text = run.text.replace(token, val_str)
                                                tokens_replaced += 1
                if slide.has_notes_slide:
                    for para in slide.notes_slide.notes_text_frame.paragraphs:
                        for run in para.runs:
                            if token in run.text:
                                run.text = run.text.replace(token, val_str)
                                tokens_replaced += 1
        result["tokens_replaced"] = tokens_replaced

    # Find/replace
    replacements_made = 0
    for rep in replacements:
        old = rep.get("old", "")
        new = rep.get("new", "")
        if not old:
            continue
        for slide in slide_list:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                replacements_made += 1
    if replacements:
        result["replacements_made"] = replacements_made

    # Table cell edits
    table_cells_edited = 0
    for te in table_edits:
        si = te.get("slide_index")
        if si is None or si < 1 or si > len(slide_list):
            continue
        slide = slide_list[si - 1]
        ti = te.get("table_index", 1)
        row = te.get("row")
        col = te.get("col")
        value = te.get("value", "")
        if row is None or col is None:
            continue
        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
                if table_count == ti:
                    tbl = shape.table
                    if 0 <= row < len(tbl.rows) and 0 <= col < len(tbl.columns):
                        tbl.cell(row, col).text = str(value)
                        table_cells_edited += 1
                    break
    if table_edits:
        result["table_cells_edited"] = table_cells_edited

    # Table cell formatting
    table_cells_formatted = 0
    for tf_entry in table_format:
        si = tf_entry.get("slide_index")
        if si is None or si < 1 or si > len(slide_list):
            continue
        slide = slide_list[si - 1]
        ti = tf_entry.get("table_index", 1)
        row = tf_entry.get("row")
        col = tf_entry.get("col")
        if row is None or col is None:
            continue
        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
                if table_count == ti:
                    tbl = shape.table
                    if 0 <= row < len(tbl.rows) and 0 <= col < len(tbl.columns):
                        cell = tbl.cell(row, col)
                        bg_color = tf_entry.get("bg_color")
                        if bg_color:
                            parsed_bg = _parse_rgb_color(bg_color)
                            if parsed_bg is not None:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = parsed_bg
                        alignment = tf_entry.get("alignment")
                        if alignment:
                            from pptx.enum.text import PP_ALIGN

                            align_map = {
                                "left": PP_ALIGN.LEFT,
                                "center": PP_ALIGN.CENTER,
                                "right": PP_ALIGN.RIGHT,
                            }
                            align_val = align_map.get(alignment.lower())
                            if align_val is not None:
                                for para in cell.text_frame.paragraphs:
                                    para.alignment = align_val
                        font_size = tf_entry.get("font_size")
                        font_bold = tf_entry.get("font_bold")
                        font_color = tf_entry.get("font_color")
                        font_name = tf_entry.get("font_name")
                        if font_size is not None or font_bold is not None or font_color or font_name:
                            parsed_fc = _parse_rgb_color(font_color) if font_color else None
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if font_size is not None:
                                        run.font.size = Pt(font_size)
                                    if font_bold is not None:
                                        run.font.bold = bool(font_bold)
                                    if parsed_fc is not None:
                                        run.font.color.rgb = parsed_fc
                                    if font_name is not None:
                                        run.font.name = str(font_name)
                        table_cells_formatted += 1
                    break
    if table_format:
        result["table_cells_formatted"] = table_cells_formatted

    # Paragraph edits — multi-paragraph text with per-paragraph formatting
    paragraphs_edited = 0
    for pe in paragraph_edits:
        si = pe.get("slide_index")
        shi = pe.get("shape_index")
        paras = pe.get("paragraphs")
        if si is None or shi is None or not paras:
            continue
        if si < 1 or si > len(slide_list):
            continue
        slide = slide_list[si - 1]
        shape_list_local = list(slide.shapes)
        if shi < 1 or shi > len(shape_list_local):
            continue
        shape = shape_list_local[shi - 1]
        if not shape.has_text_frame:
            continue
        from pptx.enum.text import PP_ALIGN

        shape.text_frame.clear()
        for p_idx, p_def in enumerate(paras):
            if p_idx == 0:
                p = shape.text_frame.paragraphs[0]
            else:
                p = shape.text_frame.add_paragraph()
            text = p_def.get("text", "")
            run = p.add_run()
            run.text = str(text)
            level = p_def.get("level")
            if level is not None:
                p.level = max(0, min(8, int(level)))
            alignment = p_def.get("alignment")
            if alignment:
                align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
                align_val = align_map.get(alignment.lower())
                if align_val is not None:
                    p.alignment = align_val
            font_size = p_def.get("font_size")
            if font_size is not None:
                run.font.size = Pt(font_size)
            font_bold = p_def.get("font_bold")
            if font_bold is not None:
                run.font.bold = bool(font_bold)
            font_color = p_def.get("font_color")
            if font_color:
                parsed = _parse_rgb_color(font_color)
                if parsed is not None:
                    run.font.color.rgb = parsed
            font_name = p_def.get("font_name")
            if font_name is not None:
                run.font.name = str(font_name)
        paragraphs_edited += 1
    if paragraph_edits:
        result["paragraphs_edited"] = paragraphs_edited

    # Placeholder edits — target by type name
    _placeholder_type_map = {
        "title": 0,
        "center_title": 3,
        "subtitle": 1,
        "body": 1,
        "slide_number": 12,
        "date": 10,
        "footer": 11,
    }
    placeholders_edited = 0
    for phe in placeholder_edits:
        si = phe.get("slide_index")
        ph_type = phe.get("placeholder_type")
        if si is None or not ph_type:
            continue
        if si < 1 or si > len(slide_list):
            continue
        slide = slide_list[si - 1]
        target_idx = _placeholder_type_map.get(ph_type.lower())
        if target_idx is None:
            continue
        target_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == target_idx:
                target_ph = ph
                break
        if target_ph is None:
            continue
        text = phe.get("text")
        if text is not None and hasattr(target_ph, "text_frame"):
            target_ph.text_frame.clear()
            target_ph.text_frame.paragraphs[0].text = str(text)
        if hasattr(target_ph, "text_frame"):
            font_size = phe.get("font_size")
            font_bold = phe.get("font_bold")
            font_color = phe.get("font_color")
            font_name = phe.get("font_name")
            if font_size is not None or font_bold is not None or font_color or font_name:
                parsed_fc = _parse_rgb_color(font_color) if font_color else None
                for para in target_ph.text_frame.paragraphs:
                    for run in para.runs:
                        if font_size is not None:
                            run.font.size = Pt(font_size)
                        if font_bold is not None:
                            run.font.bold = bool(font_bold)
                        if parsed_fc is not None:
                            run.font.color.rgb = parsed_fc
                        if font_name is not None:
                            run.font.name = str(font_name)
        placeholders_edited += 1
    if placeholder_edits:
        result["placeholders_edited"] = placeholders_edited

    # Image replacements — swap image shapes preserving position/size
    images_replaced = 0
    for ir in image_replacements:
        si = ir.get("slide_index")
        shi = ir.get("shape_index")
        img_path = ir.get("image_path")
        if si is None or shi is None or not img_path:
            continue
        if si < 1 or si > len(slide_list):
            continue
        slide = slide_list[si - 1]
        shape_list_local = list(slide.shapes)
        if shi < 1 or shi > len(shape_list_local):
            continue
        old_shape = shape_list_local[shi - 1]
        if old_shape.shape_type != 13:  # MSO_SHAPE_TYPE.PICTURE = 13
            continue
        img_resolved, img_err = validate_path(img_path, kwargs.get("working_dir", _get_working_dir()))
        if img_err or not os.path.isfile(img_resolved):
            continue
        left = old_shape.left
        top = old_shape.top
        w = old_shape.width
        h = old_shape.height
        sp_elem = old_shape._element
        sp_elem.getparent().remove(sp_elem)
        slide.shapes.add_picture(img_resolved, left, top, w, h)
        images_replaced += 1
    if image_replacements:
        result["images_replaced"] = images_replaced

    # Shape text/formatting edits
    shapes_edited = 0
    for se in shape_edits:
        si = se.get("slide_index")
        shi = se.get("shape_index")
        if si is None or shi is None or si < 1 or si > len(slide_list):
            continue
        slide = slide_list[si - 1]
        shape_list = list(slide.shapes)
        if shi < 1 or shi > len(shape_list):
            continue
        shape = shape_list[shi - 1]
        text = se.get("text")
        if text is not None and shape.has_text_frame:
            shape.text_frame.clear()
            shape.text_frame.paragraphs[0].text = str(text)
        font_size = se.get("font_size")
        font_bold = se.get("font_bold")
        font_color = se.get("font_color")
        font_name = se.get("font_name")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if font_size is not None:
                        run.font.size = Pt(font_size)
                    if font_bold is not None:
                        run.font.bold = bool(font_bold)
                    if font_color is not None:
                        rgb = _parse_rgb_color(font_color)
                        if rgb is not None:
                            run.font.color.rgb = rgb
                    if font_name is not None:
                        run.font.name = str(font_name)
        shapes_edited += 1
    if shape_edits:
        result["shapes_edited"] = shapes_edited

    # Notes edits
    notes_edited = 0
    for ne in notes_edits:
        si = ne.get("slide_index")
        text = ne.get("text")
        if si is None or text is None or si < 1 or si > len(slide_list):
            continue
        slide = slide_list[si - 1]
        if not slide.has_notes_slide:
            slide.notes_slide  # creates it
        slide.notes_slide.notes_text_frame.text = str(text)
        notes_edited += 1
    if notes_edits:
        result["notes_edited"] = notes_edited

    # Duplicate slides (before deletes)
    slides_duplicated = 0
    if duplicate_slides:
        import copy

        for di in sorted(duplicate_slides):
            if di < 1 or di > len(slide_list):
                continue
            src_slide = slide_list[di - 1]
            layout = src_slide.slide_layout
            new_slide = prs.slides.add_slide(layout)
            for elem in list(new_slide.shapes._spTree):
                if elem.tag.endswith("}sp") or elem.tag.endswith("}pic") or elem.tag.endswith("}graphicFrame"):
                    new_slide.shapes._spTree.remove(elem)
            for elem in src_slide.shapes._spTree:
                if elem.tag.endswith("}sp") or elem.tag.endswith("}pic") or elem.tag.endswith("}graphicFrame"):
                    new_slide.shapes._spTree.append(copy.deepcopy(elem))
            slides_duplicated += 1
    if duplicate_slides:
        result["slides_duplicated"] = slides_duplicated

    # Delete slides (descending order to preserve indices)
    slides_deleted = 0
    if delete_slides:
        slide_id_list = list(prs.slides._sldIdLst)
        for di in sorted(delete_slides, reverse=True):
            if di < 1 or di > len(slide_id_list):
                continue
            rel_id = slide_id_list[di - 1].get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            prs.slides._sldIdLst.remove(slide_id_list[di - 1])
            if rel_id:
                prs.part.drop_rel(rel_id)
            slides_deleted += 1
    if delete_slides:
        result["slides_deleted"] = slides_deleted

    # Append new slides
    for slide_def in slides_to_add:
        _add_slide_lib(prs, slide_def)
    if slides_to_add:
        result["slides_appended"] = len(slides_to_add)

    prs.save(resolved)
    return result


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------


def _parse_color_int(color: str) -> int:
    """Parse hex color string to COM RGB integer (BGR format)."""
    c = color.lstrip("#")
    if len(c) == 6:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return r + (g << 8) + (b << 16)
    return 0


def _parse_rgb_color(color: str) -> Optional[Any]:
    """Parse hex color string to python-pptx RGBColor.

    Returns ``None`` if the color string is invalid, letting callers decide
    how to handle the failure (typically by skipping the color application).
    """
    from pptx.dml.color import RGBColor

    c = color.lstrip("#")
    if len(c) == 3:
        c = "".join(ch * 2 for ch in c)
    try:
        return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    except (ValueError, IndexError):
        return None
