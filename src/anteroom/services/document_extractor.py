"""Text extraction from binary document formats.

Provides extract_text() which dispatches to format-specific extractors.
Dependencies (pypdf, python-docx, python-pptx, openpyxl) are optional —
extraction gracefully returns None when they are not installed.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from io import BytesIO

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class ExtractionResult:
    """Result of a text extraction attempt, with optional warnings."""

    text: str | None = None
    warnings: list[str] = field(default_factory=list)


# MIME types we can extract text from (when the library is available).
_PDF_TYPES = frozenset({"application/pdf"})
_DOCX_TYPES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }
)
_PPTX_TYPES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
)
_XLSX_TYPES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
)

EXTRACTABLE_MIME_TYPES = _PDF_TYPES | _DOCX_TYPES | _PPTX_TYPES | _XLSX_TYPES


def extract_text(data: bytes, mime_type: str) -> ExtractionResult:
    """Extract text content from binary document bytes.

    Returns an ExtractionResult with the extracted text and any warnings
    (e.g. missing library, corrupt file).
    """
    if mime_type in _PDF_TYPES:
        return _extract_pdf(data)
    if mime_type in _DOCX_TYPES:
        return _extract_docx(data)
    if mime_type in _PPTX_TYPES:
        return _extract_pptx(data)
    if mime_type in _XLSX_TYPES:
        return _extract_xlsx(data)
    return ExtractionResult()


def _extract_pdf(data: bytes) -> ExtractionResult:
    try:
        from pypdf import PdfReader
    except ImportError:
        msg = "pypdf not installed — PDF text extraction unavailable. Install with: pip install anteroom[docs]"
        logger.warning(msg)
        return ExtractionResult(warnings=[msg])
    try:
        reader = PdfReader(BytesIO(data))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        result = "\n\n".join(pages).strip()
        return ExtractionResult(text=result if result else None)
    except Exception as exc:
        msg = f"Failed to extract text from PDF: {exc}"
        logger.warning(msg, exc_info=True)
        return ExtractionResult(warnings=[msg])


def _extract_docx(data: bytes) -> ExtractionResult:
    try:
        from docx import Document
    except ImportError:
        msg = "python-docx not installed — DOCX text extraction unavailable. Install with: pip install anteroom[docs]"
        logger.warning(msg)
        return ExtractionResult(warnings=[msg])
    try:
        doc = Document(BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        result = "\n\n".join(paragraphs).strip()
        return ExtractionResult(text=result if result else None)
    except Exception as exc:
        msg = f"Failed to extract text from DOCX: {exc}"
        logger.warning(msg, exc_info=True)
        return ExtractionResult(warnings=[msg])


def _extract_pptx(data: bytes) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        msg = "python-pptx not installed — PPTX text extraction unavailable. Install with: pip install anteroom[office]"
        logger.warning(msg)
        return ExtractionResult(warnings=[msg])
    try:
        prs = Presentation(BytesIO(data))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker notes: {notes}]")
            if parts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(parts))
        result = "\n\n".join(slides).strip()
        return ExtractionResult(text=result if result else None)
    except Exception as exc:
        msg = f"Failed to extract text from PPTX: {exc}"
        logger.warning(msg, exc_info=True)
        return ExtractionResult(warnings=[msg])


def _extract_xlsx(data: bytes) -> ExtractionResult:
    try:
        from openpyxl import load_workbook
    except ImportError:
        msg = "openpyxl not installed — XLSX text extraction unavailable. Install with: pip install anteroom[office]"
        logger.warning(msg)
        return ExtractionResult(warnings=[msg])
    try:
        wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
        sheets: list[str] = []
        for ws in wb.worksheets:
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"## Sheet: {ws.title}\n" + "\n".join(rows))
        wb.close()
        result = "\n\n".join(sheets).strip()
        return ExtractionResult(text=result if result else None)
    except Exception as exc:
        msg = f"Failed to extract text from XLSX: {exc}"
        logger.warning(msg, exc_info=True)
        return ExtractionResult(warnings=[msg])
