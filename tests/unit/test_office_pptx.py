"""Tests for the PPTX tool."""

from __future__ import annotations

from unittest.mock import patch

import pytest

from anteroom.tools.office_pptx import _MAX_EDIT_OPS, _MAX_SLIDES, AVAILABLE, DEFINITION, handle, set_working_dir

_needs_pptx = pytest.mark.skipif(not AVAILABLE, reason="requires python-pptx: pip install anteroom[office]")


@pytest.fixture(autouse=True)
def _set_working_dir(tmp_path):
    set_working_dir(str(tmp_path))
    yield


class TestDefinition:
    def test_name(self):
        assert DEFINITION["name"] == "pptx"

    def test_required_params(self):
        assert "action" in DEFINITION["parameters"]["required"]
        assert "path" in DEFINITION["parameters"]["required"]


@_needs_pptx
class TestCreate:
    @pytest.mark.asyncio
    async def test_create_simple(self, tmp_path):
        result = await handle(
            action="create",
            path="test.pptx",
            slides=[
                {"title": "Slide 1", "content": "Hello world"},
                {"title": "Slide 2", "bullets": ["Point A", "Point B"]},
            ],
        )
        assert "error" not in result
        assert result["slides_created"] == 2
        assert (tmp_path / "test.pptx").exists()

    @pytest.mark.asyncio
    async def test_create_with_notes(self, tmp_path):
        result = await handle(
            action="create",
            path="notes.pptx",
            slides=[{"title": "Talk", "content": "Main point", "notes": "Speaker notes here"}],
        )
        assert "error" not in result

    @pytest.mark.asyncio
    async def test_create_no_slides(self):
        result = await handle(action="create", path="empty.pptx")
        assert "error" in result

    @pytest.mark.asyncio
    async def test_create_too_many_slides(self):
        slides = [{"title": f"Slide {i}"} for i in range(_MAX_SLIDES + 1)]
        result = await handle(action="create", path="big.pptx", slides=slides)
        assert "error" in result
        assert "Too many" in result["error"]


@_needs_pptx
class TestRead:
    @pytest.mark.asyncio
    async def test_read_simple(self, tmp_path):
        await handle(
            action="create",
            path="read_me.pptx",
            slides=[
                {"title": "Title Slide", "content": "Body text"},
            ],
        )
        result = await handle(action="read", path="read_me.pptx")
        assert "error" not in result
        assert "Title Slide" in result["content"]
        assert result["slides"] == 1

    @pytest.mark.asyncio
    async def test_read_with_notes(self, tmp_path):
        await handle(
            action="create",
            path="notes.pptx",
            slides=[{"title": "Talk", "notes": "My notes"}],
        )
        result = await handle(action="read", path="notes.pptx")
        assert "error" not in result
        assert "My notes" in result["content"]

    @pytest.mark.asyncio
    async def test_read_shows_paragraph_structure(self, tmp_path):
        """Multi-paragraph shapes should show per-paragraph detail in read output."""
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        layout = prs.slide_layouts[5]  # blank
        slide = prs.slides.add_slide(layout)
        from pptx.util import Emu

        tx_box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(5000000), Emu(5000000))
        tf = tx_box.text_frame
        # First paragraph (heading)
        tf.paragraphs[0].text = "Requirements"
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(18)
        # Second paragraph (body)
        p2 = tf.add_paragraph()
        p2.text = "Must support offline mode"
        p2.level = 1

        prs.save(str(tmp_path / "structured.pptx"))

        result = await handle(action="read", path="structured.pptx")
        assert "error" not in result
        content = result["content"]
        assert "paragraphs)" in content
        assert "P0" in content
        assert "P1" in content
        assert "Requirements" in content
        assert "bold" in content

    @pytest.mark.asyncio
    async def test_read_not_found(self):
        result = await handle(action="read", path="missing.pptx")
        assert "error" in result
        assert "not found" in result["error"].lower()

    @pytest.mark.asyncio
    async def test_read_corrupt_file(self, tmp_path):
        corrupt = tmp_path / "corrupt.pptx"
        corrupt.write_bytes(b"not a pptx file")
        result = await handle(action="read", path="corrupt.pptx")
        assert "error" in result


@_needs_pptx
class TestEdit:
    @pytest.mark.asyncio
    async def test_edit_replace(self, tmp_path):
        await handle(
            action="create",
            path="edit_me.pptx",
            slides=[{"title": "Hello world", "content": "Hello content"}],
        )
        result = await handle(
            action="edit",
            path="edit_me.pptx",
            replacements=[{"old": "Hello", "new": "Goodbye"}],
        )
        assert "error" not in result
        assert result["replacements_made"] >= 1

    @pytest.mark.asyncio
    async def test_edit_append_slides(self, tmp_path):
        await handle(
            action="create",
            path="append.pptx",
            slides=[{"title": "Original"}],
        )
        result = await handle(
            action="edit",
            path="append.pptx",
            slides=[{"title": "Appended"}],
        )
        assert "error" not in result
        assert result["slides_appended"] == 1

    @pytest.mark.asyncio
    async def test_edit_not_found(self):
        result = await handle(
            action="edit",
            path="missing.pptx",
            replacements=[{"old": "a", "new": "b"}],
        )
        assert "error" in result

    @pytest.mark.asyncio
    async def test_edit_no_operations(self, tmp_path):
        await handle(
            action="create",
            path="noop.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(action="edit", path="noop.pptx")
        assert "error" in result

    @pytest.mark.asyncio
    async def test_edit_too_many_slides(self, tmp_path):
        await handle(
            action="create",
            path="big.pptx",
            slides=[{"title": "Slide"}],
        )
        slides = [{"title": f"S{i}"} for i in range(_MAX_SLIDES + 1)]
        result = await handle(action="edit", path="big.pptx", slides=slides)
        assert "error" in result

    @pytest.mark.asyncio
    async def test_edit_exceeds_total_limit(self, tmp_path):
        # Create with many slides, then try to append more
        initial = [{"title": f"S{i}"} for i in range(_MAX_SLIDES - 1)]
        await handle(action="create", path="nearly_full.pptx", slides=initial)
        result = await handle(
            action="edit",
            path="nearly_full.pptx",
            slides=[{"title": "Extra1"}, {"title": "Extra2"}, {"title": "Extra3"}],
        )
        assert "error" in result
        assert "exceed" in result["error"].lower()

    @pytest.mark.asyncio
    async def test_edit_table_cells(self, tmp_path):
        await handle(
            action="create",
            path="table_test.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="table_test.pptx",
            slide_index=1,
            data=[["Name", "Value"], ["", ""], ["", ""]],
            rows=3,
            cols=2,
        )
        result = await handle(
            action="edit",
            path="table_test.pptx",
            table_edits=[
                {"slide_index": 1, "table_index": 1, "row": 1, "col": 0, "value": "Alice"},
                {"slide_index": 1, "table_index": 1, "row": 1, "col": 1, "value": "100"},
                {"slide_index": 1, "table_index": 1, "row": 2, "col": 0, "value": "Bob"},
                {"slide_index": 1, "table_index": 1, "row": 2, "col": 1, "value": "200"},
            ],
        )
        assert "error" not in result
        assert result["table_cells_edited"] == 4

        read_result = await handle(action="read", path="table_test.pptx")
        assert "Alice" in read_result["content"]
        assert "Bob" in read_result["content"]

    @pytest.mark.asyncio
    async def test_edit_table_cells_out_of_bounds(self, tmp_path):
        await handle(
            action="create",
            path="table_oob.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="table_oob.pptx",
            slide_index=1,
            data=[["A", "B"], ["1", "2"]],
            rows=2,
            cols=2,
        )
        result = await handle(
            action="edit",
            path="table_oob.pptx",
            table_edits=[
                {"slide_index": 1, "table_index": 1, "row": 99, "col": 0, "value": "nope"},
            ],
        )
        assert "error" not in result
        assert result["table_cells_edited"] == 0

    @pytest.mark.asyncio
    async def test_edit_shape_text(self, tmp_path):
        await handle(
            action="create",
            path="shape_edit.pptx",
            slides=[{"title": "Original Title", "content": "Original Content"}],
        )
        read_before = await handle(action="read", path="shape_edit.pptx")
        assert "Original Title" in read_before["content"]

        result = await handle(
            action="edit",
            path="shape_edit.pptx",
            shape_edits=[
                {"slide_index": 1, "shape_index": 1, "text": "New Title"},
            ],
        )
        assert "error" not in result
        assert result["shapes_edited"] == 1

    @pytest.mark.asyncio
    async def test_edit_shape_out_of_bounds(self, tmp_path):
        await handle(
            action="create",
            path="shape_oob.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(
            action="edit",
            path="shape_oob.pptx",
            shape_edits=[
                {"slide_index": 1, "shape_index": 999, "text": "nope"},
            ],
        )
        assert "error" not in result
        assert result["shapes_edited"] == 0

    @pytest.mark.asyncio
    async def test_edit_notes(self, tmp_path):
        await handle(
            action="create",
            path="notes_edit.pptx",
            slides=[{"title": "Slide 1"}, {"title": "Slide 2"}],
        )
        result = await handle(
            action="edit",
            path="notes_edit.pptx",
            notes_edits=[
                {"slide_index": 1, "text": "Speaker notes for slide 1"},
                {"slide_index": 2, "text": "Speaker notes for slide 2"},
            ],
        )
        assert "error" not in result
        assert result["notes_edited"] == 2

        read_result = await handle(action="read", path="notes_edit.pptx")
        assert "Speaker notes for slide 1" in read_result["content"]
        assert "Speaker notes for slide 2" in read_result["content"]

    @pytest.mark.asyncio
    async def test_edit_delete_slides(self, tmp_path):
        await handle(
            action="create",
            path="delete_test.pptx",
            slides=[{"title": "Keep"}, {"title": "Delete Me"}, {"title": "Also Keep"}],
        )
        result = await handle(
            action="edit",
            path="delete_test.pptx",
            delete_slides=[2],
        )
        assert "error" not in result
        assert result["slides_deleted"] == 1

        read_result = await handle(action="read", path="delete_test.pptx")
        assert read_result["slides"] == 2
        assert "Delete Me" not in read_result["content"]
        assert "Keep" in read_result["content"]

    @pytest.mark.asyncio
    async def test_edit_duplicate_slides(self, tmp_path):
        await handle(
            action="create",
            path="dup_test.pptx",
            slides=[{"title": "Original"}],
        )
        result = await handle(
            action="edit",
            path="dup_test.pptx",
            duplicate_slides=[1],
        )
        assert "error" not in result
        assert result["slides_duplicated"] == 1

        read_result = await handle(action="read", path="dup_test.pptx")
        assert read_result["slides"] == 2

    @pytest.mark.asyncio
    async def test_edit_multiple_operations(self, tmp_path):
        await handle(
            action="create",
            path="multi_edit.pptx",
            slides=[{"title": "Hello world"}],
        )
        result = await handle(
            action="edit",
            path="multi_edit.pptx",
            replacements=[{"old": "Hello", "new": "Goodbye"}],
            notes_edits=[{"slide_index": 1, "text": "New notes"}],
            slides=[{"title": "New Slide"}],
        )
        assert "error" not in result
        assert result["replacements_made"] >= 1
        assert result["notes_edited"] == 1
        assert result["slides_appended"] == 1

    @pytest.mark.asyncio
    async def test_edit_error_message_lists_all_params(self, tmp_path):
        await handle(
            action="create",
            path="err_msg.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(action="edit", path="err_msg.pptx")
        assert "error" in result
        assert "table_edits" in result["error"]
        assert "shape_edits" in result["error"]
        assert "notes_edits" in result["error"]
        assert "delete_slides" in result["error"]
        assert "duplicate_slides" in result["error"]


@_needs_pptx
class TestReadTableDisplay:
    @pytest.mark.asyncio
    async def test_read_shows_table_details(self, tmp_path):
        await handle(
            action="create",
            path="table_read.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="table_read.pptx",
            slide_index=1,
            data=[["Name", "Score"], ["Alice", "95"]],
            rows=2,
            cols=2,
        )
        result = await handle(action="read", path="table_read.pptx")
        assert "error" not in result
        assert "Table 1:" in result["content"]
        assert "Row 0:" in result["content"]
        assert "Alice" in result["content"]

    @pytest.mark.asyncio
    async def test_read_shows_shape_indices(self, tmp_path):
        await handle(
            action="create",
            path="shape_idx.pptx",
            slides=[{"title": "My Title", "content": "Body text"}],
        )
        result = await handle(action="read", path="shape_idx.pptx")
        assert "error" not in result
        assert "[Shape" in result["content"]
        assert "shapes)" in result["content"]


@_needs_pptx
class TestPathValidation:
    @pytest.mark.asyncio
    async def test_blocked_system_path_rejected(self):
        result = await handle(
            action="create",
            path="/etc/shadow",
            slides=[{"title": "x"}],
        )
        assert "error" in result

    @pytest.mark.asyncio
    async def test_null_bytes_rejected(self):
        result = await handle(action="read", path="test\x00.pptx")
        assert "error" in result


@_needs_pptx
class TestUnknownAction:
    @pytest.mark.asyncio
    async def test_unknown_action(self):
        result = await handle(action="delete", path="test.pptx")
        assert "error" in result
        assert "Unknown action" in result["error"]


class TestGracefulDegradation:
    @pytest.mark.asyncio
    async def test_unavailable(self):
        with patch("anteroom.tools.office_pptx.AVAILABLE", False):
            result = await handle(action="read", path="test.pptx")
            assert "error" in result
            assert "pip install" in result["error"]


# ---------------------------------------------------------------------------
# New action tests (16 actions added for COM office backend)
# ---------------------------------------------------------------------------


@_needs_pptx
class TestInsertImage:
    async def _create_pptx(self, name: str) -> dict:
        return await handle(
            action="create",
            path=name,
            slides=[{"title": "Slide 1", "content": "Body"}],
        )

    @pytest.mark.asyncio
    async def test_insert_image_success(self, tmp_path):
        await self._create_pptx("img.pptx")
        # Create a minimal valid PNG (1x1 pixel)
        import struct
        import zlib

        def _make_png() -> bytes:
            sig = b"\x89PNG\r\n\x1a\n"

            def _chunk(ctype: bytes, data: bytes) -> bytes:
                c = ctype + data
                return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

            ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
            raw = zlib.compress(b"\x00\xff\x00\x00")
            return sig + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", raw) + _chunk(b"IEND", b"")

        img_file = tmp_path / "test.png"
        img_file.write_bytes(_make_png())

        result = await handle(
            action="insert_image",
            path="img.pptx",
            slide_index=1,
            image_path="test.png",
            left=1,
            top=1,
            width=4,
            height=3,
        )
        assert "error" not in result
        assert "Inserted image" in result["result"]
        assert result["path"] == "img.pptx"

    @pytest.mark.asyncio
    async def test_insert_image_missing_image_path(self, tmp_path):
        await self._create_pptx("img2.pptx")
        result = await handle(action="insert_image", path="img2.pptx", slide_index=1)
        assert "error" in result
        assert "image_path" in result["error"]

    @pytest.mark.asyncio
    async def test_insert_image_file_not_found(self, tmp_path):
        await self._create_pptx("img3.pptx")
        result = await handle(
            action="insert_image",
            path="img3.pptx",
            slide_index=1,
            image_path="nonexistent.png",
        )
        assert "error" in result
        assert "not found" in result["error"].lower()

    @pytest.mark.asyncio
    async def test_insert_image_pptx_not_found(self):
        result = await handle(
            action="insert_image",
            path="missing.pptx",
            slide_index=1,
            image_path="test.png",
        )
        assert "error" in result

    @pytest.mark.asyncio
    async def test_insert_image_missing_slide_index(self, tmp_path):
        await self._create_pptx("img4.pptx")
        img_file = tmp_path / "tiny.png"
        # Write a minimal PNG
        import struct
        import zlib

        sig = b"\x89PNG\r\n\x1a\n"

        def _chunk(ctype: bytes, data: bytes) -> bytes:
            c = ctype + data
            return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

        ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
        raw = zlib.compress(b"\x00\xff\x00\x00")
        img_file.write_bytes(sig + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", raw) + _chunk(b"IEND", b""))

        result = await handle(
            action="insert_image",
            path="img4.pptx",
            image_path="tiny.png",
        )
        assert "error" in result
        assert "slide_index" in result["error"]

    @pytest.mark.asyncio
    async def test_insert_image_slide_index_out_of_range(self, tmp_path):
        await self._create_pptx("img5.pptx")
        img_file = tmp_path / "tiny2.png"
        import struct
        import zlib

        sig = b"\x89PNG\r\n\x1a\n"

        def _chunk(ctype: bytes, data: bytes) -> bytes:
            c = ctype + data
            return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

        ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
        raw = zlib.compress(b"\x00\xff\x00\x00")
        img_file.write_bytes(sig + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", raw) + _chunk(b"IEND", b""))

        result = await handle(
            action="insert_image",
            path="img5.pptx",
            slide_index=99,
            image_path="tiny2.png",
        )
        assert "error" in result
        assert "out of range" in result["error"]


@_needs_pptx
class TestInsertShape:
    async def _create_pptx(self, name: str) -> dict:
        return await handle(
            action="create",
            path=name,
            slides=[{"title": "Slide 1"}],
        )

    @pytest.mark.asyncio
    async def test_insert_rectangle(self, tmp_path):
        await self._create_pptx("shape.pptx")
        result = await handle(
            action="insert_shape",
            path="shape.pptx",
            slide_index=1,
            shape_type="rectangle",
            left=2,
            top=2,
            width=3,
            height=2,
        )
        assert "error" not in result
        assert "Inserted rectangle" in result["result"]

    @pytest.mark.asyncio
    async def test_insert_oval(self, tmp_path):
        await self._create_pptx("shape_oval.pptx")
        result = await handle(
            action="insert_shape",
            path="shape_oval.pptx",
            slide_index=1,
            shape_type="oval",
        )
        assert "error" not in result

    @pytest.mark.asyncio
    async def test_insert_all_shape_types(self, tmp_path):
        shape_types = [
            "rectangle",
            "oval",
            "triangle",
            "right_arrow",
            "left_arrow",
            "diamond",
            "rounded_rectangle",
            "star",
        ]
        for shape_type in shape_types:
            await self._create_pptx(f"shape_{shape_type}.pptx")
            result = await handle(
                action="insert_shape",
                path=f"shape_{shape_type}.pptx",
                slide_index=1,
                shape_type=shape_type,
            )
            assert "error" not in result, f"Failed for shape_type={shape_type}: {result}"

    @pytest.mark.asyncio
    async def test_insert_shape_unknown_type(self, tmp_path):
        await self._create_pptx("shape_bad.pptx")
        result = await handle(
            action="insert_shape",
            path="shape_bad.pptx",
            slide_index=1,
            shape_type="hexagon_of_doom",
        )
        assert "error" in result
        assert "Unknown shape_type" in result["error"]

    @pytest.mark.asyncio
    async def test_insert_shape_missing_type(self, tmp_path):
        await self._create_pptx("shape_notype.pptx")
        result = await handle(
            action="insert_shape",
            path="shape_notype.pptx",
            slide_index=1,
        )
        assert "error" in result
        assert "shape_type" in result["error"]

    @pytest.mark.asyncio
    async def test_insert_shape_missing_slide_index(self, tmp_path):
        await self._create_pptx("shape_nosi.pptx")
        result = await handle(
            action="insert_shape",
            path="shape_nosi.pptx",
            shape_type="rectangle",
        )
        assert "error" in result
        assert "slide_index" in result["error"]

    @pytest.mark.asyncio
    async def test_insert_shape_slide_out_of_range(self, tmp_path):
        await self._create_pptx("shape_oor.pptx")
        result = await handle(
            action="insert_shape",
            path="shape_oor.pptx",
            slide_index=99,
            shape_type="rectangle",
        )
        assert "error" in result
        assert "out of range" in result["error"]

    @pytest.mark.asyncio
    async def test_insert_shape_file_not_found(self):
        result = await handle(
            action="insert_shape",
            path="missing.pptx",
            slide_index=1,
            shape_type="rectangle",
        )
        assert "error" in result


@_needs_pptx
class TestFormatShape:
    async def _create_with_shape(self, name: str) -> dict:
        await handle(action="create", path=name, slides=[{"title": "Slide 1"}])
        await handle(
            action="insert_shape",
            path=name,
            slide_index=1,
            shape_type="rectangle",
            left=2,
            top=2,
            width=3,
            height=2,
        )
        return await handle(action="read", path=name)

    @pytest.mark.asyncio
    async def test_format_shape_fill_color(self, tmp_path):
        await self._create_with_shape("fmt.pptx")
        # The shape we inserted is not at index 1 (title placeholder is),
        # so we need to find its index. With the lib backend, shapes are
        # ordered: title placeholder, content placeholder (if any), then added shapes.
        # For a layout-1 slide with title+content, our added shape is index 3.
        # We can just try a reasonable index; the important thing is no crash on valid format.
        from pptx import Presentation

        prs = Presentation(str(tmp_path / "fmt.pptx"))
        num_shapes = len(list(list(prs.slides)[0].shapes))

        result = await handle(
            action="format_shape",
            path="fmt.pptx",
            slide_index=1,
            shape_index=num_shapes,
            format={"fill_color": "#FF0000"},
        )
        assert "error" not in result
        assert "Formatted shape" in result["result"]

    @pytest.mark.asyncio
    async def test_format_shape_line_color_and_width(self, tmp_path):
        await self._create_with_shape("fmt_line.pptx")
        from pptx import Presentation

        prs = Presentation(str(tmp_path / "fmt_line.pptx"))
        num_shapes = len(list(list(prs.slides)[0].shapes))

        result = await handle(
            action="format_shape",
            path="fmt_line.pptx",
            slide_index=1,
            shape_index=num_shapes,
            format={"line_color": "#00FF00", "line_width": 2},
        )
        assert "error" not in result

    @pytest.mark.asyncio
    async def test_format_shape_text_and_font(self, tmp_path):
        await self._create_with_shape("fmt_text.pptx")
        from pptx import Presentation

        prs = Presentation(str(tmp_path / "fmt_text.pptx"))
        num_shapes = len(list(list(prs.slides)[0].shapes))

        result = await handle(
            action="format_shape",
            path="fmt_text.pptx",
            slide_index=1,
            shape_index=num_shapes,
            format={"text": "Hello", "font_size": 14, "font_bold": True, "font_color": "#0000FF"},
        )
        assert "error" not in result

    @pytest.mark.asyncio
    async def test_format_shape_missing_format(self, tmp_path):
        await self._create_with_shape("fmt_nofmt.pptx")
        result = await handle(
            action="format_shape",
            path="fmt_nofmt.pptx",
            slide_index=1,
            shape_index=1,
        )
        assert "error" in result
        assert "format" in result["error"].lower()

    @pytest.mark.asyncio
    async def test_format_shape_missing_slide_index(self, tmp_path):
        await self._create_with_shape("fmt_nosi.pptx")
        result = await handle(
            action="format_shape",
            path="fmt_nosi.pptx",
            shape_index=1,
            format={"fill_color": "#FF0000"},
        )
        assert "error" in result
        assert "slide_index" in result["error"]

    @pytest.mark.asyncio
    async def test_format_shape_missing_shape_index(self, tmp_path):
        await self._create_with_shape("fmt_nosh.pptx")
        result = await handle(
            action="format_shape",
            path="fmt_nosh.pptx",
            slide_index=1,
            format={"fill_color": "#FF0000"},
        )
        assert "error" in result
        assert "shape_index" in result["error"]

    @pytest.mark.asyncio
    async def test_format_shape_index_out_of_range(self, tmp_path):
        await self._create_with_shape("fmt_oor.pptx")
        result = await handle(
            action="format_shape",
            path="fmt_oor.pptx",
            slide_index=1,
            shape_index=99,
            format={"fill_color": "#FF0000"},
        )
        assert "error" in result
        assert "out of range" in result["error"]

    @pytest.mark.asyncio
    async def test_format_shape_file_not_found(self):
        result = await handle(
            action="format_shape",
            path="missing.pptx",
            slide_index=1,
            shape_index=1,
            format={"fill_color": "#FF0000"},
        )
        assert "error" in result


@_needs_pptx
class TestMasterLayout:
    async def _create_pptx(self, name: str) -> dict:
        return await handle(
            action="create",
            path=name,
            slides=[{"title": "Slide 1"}],
        )

    @pytest.mark.asyncio
    async def test_master_layout_list(self, tmp_path):
        await self._create_pptx("layout.pptx")
        result = await handle(
            action="master_layout",
            path="layout.pptx",
            operation="list",
        )
        assert "error" not in result
        assert "layouts" in result
        assert "count" in result
        assert isinstance(result["layouts"], list)
        assert result["count"] > 0

    @pytest.mark.asyncio
    async def test_master_layout_list_default_operation(self, tmp_path):
        await self._create_pptx("layout_default.pptx")
        result = await handle(
            action="master_layout",
            path="layout_default.pptx",
        )
        assert "error" not in result
        assert "layouts" in result

    @pytest.mark.asyncio
    async def test_master_layout_apply_com_only(self, tmp_path):
        await self._create_pptx("layout_apply.pptx")
        result = await handle(
            action="master_layout",
            path="layout_apply.pptx",
            operation="apply",
            slide_index=1,
            layout_index=0,
        )
        assert "error" in result
        assert "COM" in result["error"]

    @pytest.mark.asyncio
    async def test_master_layout_apply_missing_params(self, tmp_path):
        await self._create_pptx("layout_nop.pptx")
        result = await handle(
            action="master_layout",
            path="layout_nop.pptx",
            operation="apply",
        )
        assert "error" in result
        assert "COM" in result["error"]

    @pytest.mark.asyncio
    async def test_master_layout_unknown_operation(self, tmp_path):
        await self._create_pptx("layout_bad.pptx")
        result = await handle(
            action="master_layout",
            path="layout_bad.pptx",
            operation="frobnicate",
        )
        assert "error" in result
        assert "Unknown operation" in result["error"]

    @pytest.mark.asyncio
    async def test_master_layout_file_not_found(self):
        result = await handle(
            action="master_layout",
            path="missing.pptx",
            operation="list",
        )
        assert "error" in result

    @pytest.mark.asyncio
    async def test_master_layout_apply_layout_out_of_range(self, tmp_path):
        await self._create_pptx("layout_oor.pptx")
        result = await handle(
            action="master_layout",
            path="layout_oor.pptx",
            operation="apply",
            slide_index=1,
            layout_index=999,
        )
        assert "error" in result
        assert "COM" in result["error"]


@_needs_pptx
class TestReorderSlides:
    async def _create_multi(self, name: str) -> dict:
        return await handle(
            action="create",
            path=name,
            slides=[
                {"title": "Slide A"},
                {"title": "Slide B"},
                {"title": "Slide C"},
            ],
        )

    @pytest.mark.asyncio
    async def test_reorder_delete(self, tmp_path):
        await self._create_multi("reorder.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder.pptx",
            operation="delete",
            slide_index=2,
        )
        assert "error" not in result
        assert "Deleted slide 2" in result["result"]

        read_result = await handle(action="read", path="reorder.pptx")
        assert read_result["slides"] == 2

    @pytest.mark.asyncio
    async def test_reorder_duplicate(self, tmp_path):
        await self._create_multi("reorder_dup.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder_dup.pptx",
            operation="duplicate",
            slide_index=1,
        )
        assert "error" in result
        assert "COM backend" in result["error"]

    @pytest.mark.asyncio
    async def test_reorder_move(self, tmp_path):
        await self._create_multi("reorder_mv.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder_mv.pptx",
            operation="move",
            slide_index=1,
            new_position=3,
        )
        assert "error" not in result
        assert "Moved slide 1 to position 3" in result["result"]

    @pytest.mark.asyncio
    async def test_reorder_move_missing_new_position(self, tmp_path):
        await self._create_multi("reorder_nonp.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder_nonp.pptx",
            operation="move",
            slide_index=1,
        )
        assert "error" in result
        assert "new_position" in result["error"]

    @pytest.mark.asyncio
    async def test_reorder_move_out_of_range(self, tmp_path):
        await self._create_multi("reorder_oor.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder_oor.pptx",
            operation="move",
            slide_index=1,
            new_position=99,
        )
        assert "error" in result
        assert "out of range" in result["error"]

    @pytest.mark.asyncio
    async def test_reorder_missing_slide_index(self, tmp_path):
        await self._create_multi("reorder_nosi.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder_nosi.pptx",
            operation="delete",
        )
        assert "error" in result
        assert "slide_index" in result["error"]

    @pytest.mark.asyncio
    async def test_reorder_slide_index_out_of_range(self, tmp_path):
        await self._create_multi("reorder_sioor.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder_sioor.pptx",
            operation="delete",
            slide_index=99,
        )
        assert "error" in result
        assert "out of range" in result["error"]

    @pytest.mark.asyncio
    async def test_reorder_unknown_operation(self, tmp_path):
        await self._create_multi("reorder_bad.pptx")
        result = await handle(
            action="reorder_slides",
            path="reorder_bad.pptx",
            operation="shuffle",
            slide_index=1,
        )
        assert "error" in result
        assert "Unknown operation" in result["error"]

    @pytest.mark.asyncio
    async def test_reorder_file_not_found(self):
        result = await handle(
            action="reorder_slides",
            path="missing.pptx",
            operation="delete",
            slide_index=1,
        )
        assert "error" in result


@_needs_pptx
class TestEmbedTable:
    async def _create_pptx(self, name: str) -> dict:
        return await handle(
            action="create",
            path=name,
            slides=[{"title": "Slide 1"}],
        )

    @pytest.mark.asyncio
    async def test_embed_table_with_data(self, tmp_path):
        await self._create_pptx("table.pptx")
        result = await handle(
            action="embed_table",
            path="table.pptx",
            slide_index=1,
            data=[
                ["Name", "Age", "City"],
                ["Alice", "30", "NYC"],
                ["Bob", "25", "LA"],
            ],
        )
        assert "error" not in result
        assert "3x3 table" in result["result"]
        assert result["path"] == "table.pptx"

    @pytest.mark.asyncio
    async def test_embed_table_with_rows_cols(self, tmp_path):
        await self._create_pptx("table_rc.pptx")
        result = await handle(
            action="embed_table",
            path="table_rc.pptx",
            slide_index=1,
            rows=4,
            cols=2,
        )
        assert "error" not in result
        assert "4x2 table" in result["result"]

    @pytest.mark.asyncio
    async def test_embed_table_missing_dimensions(self, tmp_path):
        await self._create_pptx("table_nodim.pptx")
        result = await handle(
            action="embed_table",
            path="table_nodim.pptx",
            slide_index=1,
        )
        assert "error" in result
        assert "rows and cols" in result["error"]

    @pytest.mark.asyncio
    async def test_embed_table_missing_slide_index(self, tmp_path):
        await self._create_pptx("table_nosi.pptx")
        result = await handle(
            action="embed_table",
            path="table_nosi.pptx",
            rows=2,
            cols=2,
        )
        assert "error" in result
        assert "slide_index" in result["error"]

    @pytest.mark.asyncio
    async def test_embed_table_slide_out_of_range(self, tmp_path):
        await self._create_pptx("table_oor.pptx")
        result = await handle(
            action="embed_table",
            path="table_oor.pptx",
            slide_index=99,
            rows=2,
            cols=2,
        )
        assert "error" in result
        assert "out of range" in result["error"]

    @pytest.mark.asyncio
    async def test_embed_table_file_not_found(self):
        result = await handle(
            action="embed_table",
            path="missing.pptx",
            slide_index=1,
            rows=2,
            cols=2,
        )
        assert "error" in result

    @pytest.mark.asyncio
    async def test_embed_table_with_custom_position(self, tmp_path):
        await self._create_pptx("table_pos.pptx")
        result = await handle(
            action="embed_table",
            path="table_pos.pptx",
            slide_index=1,
            rows=2,
            cols=3,
            left=0.5,
            top=1.5,
            width=9,
            height=4,
        )
        assert "error" not in result


@_needs_pptx
class TestHyperlinks:
    async def _create_pptx(self, name: str) -> dict:
        return await handle(
            action="create",
            path=name,
            slides=[{"title": "Click Here", "content": "Some body text"}],
        )

    @pytest.mark.asyncio
    async def test_hyperlinks_add(self, tmp_path):
        await self._create_pptx("link.pptx")
        result = await handle(
            action="hyperlinks",
            path="link.pptx",
            operation="add",
            slide_index=1,
            shape_index=1,
            url="https://example.com",
            display_text="Example",
        )
        assert "error" not in result
        assert "Added hyperlink" in result["result"]

    @pytest.mark.asyncio
    async def test_hyperlinks_add_without_display_text(self, tmp_path):
        await self._create_pptx("link_nodt.pptx")
        result = await handle(
            action="hyperlinks",
            path="link_nodt.pptx",
            operation="add",
            slide_index=1,
            shape_index=1,
            url="https://example.com",
        )
        assert "error" not in result

    @pytest.mark.asyncio
    async def test_hyperlinks_list(self, tmp_path):
        await self._create_pptx("link_list.pptx")
        # Add a hyperlink first
        await handle(
            action="hyperlinks",
            path="link_list.pptx",
            operation="add",
            slide_index=1,
            shape_index=1,
            url="https://example.com",
            display_text="Example",
        )
        result = await handle(
            action="hyperlinks",
            path="link_list.pptx",
            operation="list",
            slide_index=1,
        )
        assert "error" not in result
        assert "hyperlinks" in result
        assert isinstance(result["hyperlinks"], list)
        assert result["slide_index"] == 1

    @pytest.mark.asyncio
    async def test_hyperlinks_add_missing_url(self, tmp_path):
        await self._create_pptx("link_nourl.pptx")
        result = await handle(
            action="hyperlinks",
            path="link_nourl.pptx",
            operation="add",
            slide_index=1,
            shape_index=1,
        )
        assert "error" in result
        assert "url" in result["error"]

    @pytest.mark.asyncio
    async def test_hyperlinks_add_missing_shape_index(self, tmp_path):
        await self._create_pptx("link_nosh.pptx")
        result = await handle(
            action="hyperlinks",
            path="link_nosh.pptx",
            operation="add",
            slide_index=1,
            url="https://example.com",
        )
        assert "error" in result
        assert "shape_index" in result["error"]

    @pytest.mark.asyncio
    async def test_hyperlinks_missing_slide_index(self, tmp_path):
        await self._create_pptx("link_nosi.pptx")
        result = await handle(
            action="hyperlinks",
            path="link_nosi.pptx",
            operation="list",
        )
        assert "error" in result
        assert "slide_index" in result["error"]

    @pytest.mark.asyncio
    async def test_hyperlinks_shape_index_out_of_range(self, tmp_path):
        await self._create_pptx("link_oor.pptx")
        result = await handle(
            action="hyperlinks",
            path="link_oor.pptx",
            operation="add",
            slide_index=1,
            shape_index=99,
            url="https://example.com",
        )
        assert "error" in result
        assert "out of range" in result["error"]

    @pytest.mark.asyncio
    async def test_hyperlinks_file_not_found(self):
        result = await handle(
            action="hyperlinks",
            path="missing.pptx",
            slide_index=1,
            operation="list",
        )
        assert "error" in result


@_needs_pptx
class TestHeadersFooters:
    async def _create_pptx(self, name: str) -> dict:
        return await handle(
            action="create",
            path=name,
            slides=[{"title": "Slide 1"}, {"title": "Slide 2"}],
        )

    @pytest.mark.asyncio
    async def test_headers_footers_list(self, tmp_path):
        await self._create_pptx("hf.pptx")
        result = await handle(
            action="headers_footers",
            path="hf.pptx",
            operation="list",
        )
        assert "error" not in result
        assert "headers_footers" in result

    @pytest.mark.asyncio
    async def test_headers_footers_set_footer(self, tmp_path):
        await self._create_pptx("hf_set.pptx")
        result = await handle(
            action="headers_footers",
            path="hf_set.pptx",
            operation="set",
            footer_text="My Footer",
        )
        assert "error" not in result
        assert "Updated headers/footers" in result["result"]

    @pytest.mark.asyncio
    async def test_headers_footers_set_on_specific_slide(self, tmp_path):
        await self._create_pptx("hf_slide.pptx")
        result = await handle(
            action="headers_footers",
            path="hf_slide.pptx",
            operation="set",
            slide_index=1,
            footer_text="Slide 1 Footer",
        )
        assert "error" not in result
        assert "slide 1" in result["result"]

    @pytest.mark.asyncio
    async def test_headers_footers_set_slide_numbers(self, tmp_path):
        await self._create_pptx("hf_num.pptx")
        result = await handle(
            action="headers_footers",
            path="hf_num.pptx",
            operation="set",
            slide_numbers=True,
        )
        assert "error" not in result

    @pytest.mark.asyncio
    async def test_headers_footers_set_date_time(self, tmp_path):
        await self._create_pptx("hf_dt.pptx")
        result = await handle(
            action="headers_footers",
            path="hf_dt.pptx",
            operation="set",
            date_time="2026-01-01",
        )
        assert "error" not in result

    @pytest.mark.asyncio
    async def test_headers_footers_default_operation_is_set(self, tmp_path):
        await self._create_pptx("hf_default.pptx")
        result = await handle(
            action="headers_footers",
            path="hf_default.pptx",
            footer_text="Default op footer",
        )
        assert "error" not in result
        assert "Updated headers/footers" in result["result"]

    @pytest.mark.asyncio
    async def test_headers_footers_slide_out_of_range(self, tmp_path):
        await self._create_pptx("hf_oor.pptx")
        result = await handle(
            action="headers_footers",
            path="hf_oor.pptx",
            operation="set",
            slide_index=99,
            footer_text="Oops",
        )
        assert "error" in result
        assert "out of range" in result["error"]

    @pytest.mark.asyncio
    async def test_headers_footers_file_not_found(self):
        result = await handle(
            action="headers_footers",
            path="missing.pptx",
            operation="list",
        )
        assert "error" in result


# ---------------------------------------------------------------------------
# New edit features: template_fill, table_format, paragraph_edits,
# placeholder_edits, image_replacements
# ---------------------------------------------------------------------------


@_needs_pptx
class TestTemplateFill:
    @pytest.mark.asyncio
    async def test_template_fill_in_text(self, tmp_path):
        await handle(
            action="create",
            path="tmpl.pptx",
            slides=[{"title": "Hello {{name}}", "content": "Date: {{date}}"}],
        )
        result = await handle(
            action="edit",
            path="tmpl.pptx",
            template_fill={"name": "World", "date": "2025-01-01"},
        )
        assert "error" not in result
        assert result["tokens_replaced"] >= 2

        read_result = await handle(action="read", path="tmpl.pptx")
        assert "World" in read_result["content"]
        assert "2025-01-01" in read_result["content"]
        assert "{{name}}" not in read_result["content"]

    @pytest.mark.asyncio
    async def test_template_fill_in_table(self, tmp_path):
        await handle(
            action="create",
            path="tmpl_tbl.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="tmpl_tbl.pptx",
            slide_index=1,
            data=[["Name", "Score"], ["{{player}}", "{{score}}"]],
            rows=2,
            cols=2,
        )
        result = await handle(
            action="edit",
            path="tmpl_tbl.pptx",
            template_fill={"player": "Alice", "score": "100"},
        )
        assert "error" not in result
        assert result["tokens_replaced"] >= 2

        read_result = await handle(action="read", path="tmpl_tbl.pptx")
        assert "Alice" in read_result["content"]
        assert "100" in read_result["content"]

    @pytest.mark.asyncio
    async def test_template_fill_in_notes(self, tmp_path):
        await handle(
            action="create",
            path="tmpl_notes.pptx",
            slides=[{"title": "Slide", "notes": "Speaker: {{speaker}}"}],
        )
        result = await handle(
            action="edit",
            path="tmpl_notes.pptx",
            template_fill={"speaker": "Bob"},
        )
        assert "error" not in result
        assert result["tokens_replaced"] >= 1

        read_result = await handle(action="read", path="tmpl_notes.pptx")
        assert "Bob" in read_result["content"]

    @pytest.mark.asyncio
    async def test_template_fill_no_match(self, tmp_path):
        await handle(
            action="create",
            path="tmpl_no.pptx",
            slides=[{"title": "Plain title"}],
        )
        result = await handle(
            action="edit",
            path="tmpl_no.pptx",
            template_fill={"missing": "value"},
        )
        assert "error" not in result
        assert result["tokens_replaced"] == 0


@_needs_pptx
class TestTableFormat:
    @pytest.mark.asyncio
    async def test_table_format_bg_color(self, tmp_path):
        await handle(
            action="create",
            path="tf.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="tf.pptx",
            slide_index=1,
            data=[["A", "B"], ["1", "2"]],
            rows=2,
            cols=2,
        )
        result = await handle(
            action="edit",
            path="tf.pptx",
            table_format=[
                {"slide_index": 1, "table_index": 1, "row": 0, "col": 0, "bg_color": "#FF0000"},
            ],
        )
        assert "error" not in result
        assert result["table_cells_formatted"] == 1

    @pytest.mark.asyncio
    async def test_table_format_font_styling(self, tmp_path):
        await handle(
            action="create",
            path="tf_font.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="tf_font.pptx",
            slide_index=1,
            data=[["Header", "Value"], ["X", "Y"]],
            rows=2,
            cols=2,
        )
        result = await handle(
            action="edit",
            path="tf_font.pptx",
            table_format=[
                {
                    "slide_index": 1,
                    "table_index": 1,
                    "row": 0,
                    "col": 0,
                    "font_size": 14,
                    "font_bold": True,
                    "font_color": "#0000FF",
                    "font_name": "Arial",
                },
            ],
        )
        assert "error" not in result
        assert result["table_cells_formatted"] == 1

    @pytest.mark.asyncio
    async def test_table_format_alignment(self, tmp_path):
        await handle(
            action="create",
            path="tf_align.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="tf_align.pptx",
            slide_index=1,
            data=[["Left", "Center"], ["A", "B"]],
            rows=2,
            cols=2,
        )
        result = await handle(
            action="edit",
            path="tf_align.pptx",
            table_format=[
                {"slide_index": 1, "table_index": 1, "row": 0, "col": 1, "alignment": "center"},
            ],
        )
        assert "error" not in result
        assert result["table_cells_formatted"] == 1

    @pytest.mark.asyncio
    async def test_table_format_out_of_bounds(self, tmp_path):
        await handle(
            action="create",
            path="tf_oob.pptx",
            slides=[{"title": "Data"}],
        )
        await handle(
            action="embed_table",
            path="tf_oob.pptx",
            slide_index=1,
            data=[["A"]],
            rows=1,
            cols=1,
        )
        result = await handle(
            action="edit",
            path="tf_oob.pptx",
            table_format=[
                {"slide_index": 1, "table_index": 1, "row": 99, "col": 0, "bg_color": "#FF0000"},
            ],
        )
        assert "error" not in result
        assert result["table_cells_formatted"] == 0


@_needs_pptx
class TestParagraphEdits:
    @pytest.mark.asyncio
    async def test_paragraph_edits_multi_para(self, tmp_path):
        await handle(
            action="create",
            path="para.pptx",
            slides=[{"title": "Title", "content": "Original"}],
        )
        result = await handle(
            action="edit",
            path="para.pptx",
            paragraph_edits=[
                {
                    "slide_index": 1,
                    "shape_index": 2,
                    "paragraphs": [
                        {"text": "First paragraph"},
                        {"text": "Second paragraph"},
                        {"text": "Third paragraph"},
                    ],
                }
            ],
        )
        assert "error" not in result
        assert result["paragraphs_edited"] == 1

        read_result = await handle(action="read", path="para.pptx")
        assert "First paragraph" in read_result["content"]
        assert "Third paragraph" in read_result["content"]

    @pytest.mark.asyncio
    async def test_paragraph_edits_with_formatting(self, tmp_path):
        await handle(
            action="create",
            path="para_fmt.pptx",
            slides=[{"title": "Title", "content": "Original"}],
        )
        result = await handle(
            action="edit",
            path="para_fmt.pptx",
            paragraph_edits=[
                {
                    "slide_index": 1,
                    "shape_index": 2,
                    "paragraphs": [
                        {"text": "Bold heading", "font_bold": True, "font_size": 24},
                        {"text": "Normal text", "alignment": "center"},
                    ],
                }
            ],
        )
        assert "error" not in result
        assert result["paragraphs_edited"] == 1

    @pytest.mark.asyncio
    async def test_paragraph_edits_with_levels(self, tmp_path):
        await handle(
            action="create",
            path="para_lvl.pptx",
            slides=[{"title": "Title", "content": "Original"}],
        )
        result = await handle(
            action="edit",
            path="para_lvl.pptx",
            paragraph_edits=[
                {
                    "slide_index": 1,
                    "shape_index": 2,
                    "paragraphs": [
                        {"text": "Top level", "level": 0},
                        {"text": "Indented", "level": 1},
                        {"text": "More indented", "level": 2},
                    ],
                }
            ],
        )
        assert "error" not in result
        assert result["paragraphs_edited"] == 1

    @pytest.mark.asyncio
    async def test_paragraph_edits_shape_out_of_bounds(self, tmp_path):
        await handle(
            action="create",
            path="para_oob.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(
            action="edit",
            path="para_oob.pptx",
            paragraph_edits=[
                {
                    "slide_index": 1,
                    "shape_index": 999,
                    "paragraphs": [{"text": "nope"}],
                }
            ],
        )
        assert "error" not in result
        assert result["paragraphs_edited"] == 0

    @pytest.mark.asyncio
    async def test_paragraph_edits_missing_paragraphs(self, tmp_path):
        await handle(
            action="create",
            path="para_empty.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(
            action="edit",
            path="para_empty.pptx",
            paragraph_edits=[
                {"slide_index": 1, "shape_index": 1},
            ],
        )
        assert "error" not in result
        assert result["paragraphs_edited"] == 0


@_needs_pptx
class TestPlaceholderEdits:
    @pytest.mark.asyncio
    async def test_placeholder_edit_title(self, tmp_path):
        await handle(
            action="create",
            path="ph.pptx",
            slides=[{"title": "Old Title", "content": "Body"}],
        )
        result = await handle(
            action="edit",
            path="ph.pptx",
            placeholder_edits=[
                {"slide_index": 1, "placeholder_type": "title", "text": "New Title"},
            ],
        )
        assert "error" not in result
        assert result["placeholders_edited"] == 1

        read_result = await handle(action="read", path="ph.pptx")
        assert "New Title" in read_result["content"]

    @pytest.mark.asyncio
    async def test_placeholder_edit_body(self, tmp_path):
        await handle(
            action="create",
            path="ph_body.pptx",
            slides=[{"title": "Title", "content": "Old body"}],
        )
        result = await handle(
            action="edit",
            path="ph_body.pptx",
            placeholder_edits=[
                {"slide_index": 1, "placeholder_type": "body", "text": "New body text"},
            ],
        )
        assert "error" not in result
        assert result["placeholders_edited"] == 1

        read_result = await handle(action="read", path="ph_body.pptx")
        assert "New body text" in read_result["content"]

    @pytest.mark.asyncio
    async def test_placeholder_edit_with_formatting(self, tmp_path):
        await handle(
            action="create",
            path="ph_fmt.pptx",
            slides=[{"title": "Title"}],
        )
        result = await handle(
            action="edit",
            path="ph_fmt.pptx",
            placeholder_edits=[
                {
                    "slide_index": 1,
                    "placeholder_type": "title",
                    "text": "Styled Title",
                    "font_size": 36,
                    "font_bold": True,
                },
            ],
        )
        assert "error" not in result
        assert result["placeholders_edited"] == 1

    @pytest.mark.asyncio
    async def test_placeholder_edit_unknown_type(self, tmp_path):
        await handle(
            action="create",
            path="ph_unk.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(
            action="edit",
            path="ph_unk.pptx",
            placeholder_edits=[
                {"slide_index": 1, "placeholder_type": "nonexistent", "text": "nope"},
            ],
        )
        assert "error" not in result
        assert result["placeholders_edited"] == 0

    @pytest.mark.asyncio
    async def test_placeholder_edit_slide_out_of_bounds(self, tmp_path):
        await handle(
            action="create",
            path="ph_oob.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(
            action="edit",
            path="ph_oob.pptx",
            placeholder_edits=[
                {"slide_index": 99, "placeholder_type": "title", "text": "nope"},
            ],
        )
        assert "error" not in result
        assert result["placeholders_edited"] == 0


@_needs_pptx
class TestImageReplacements:
    @pytest.mark.asyncio
    async def test_image_replacement_basic(self, tmp_path):
        # Create a minimal 1x1 PNG
        import struct
        import zlib

        def _make_png() -> bytes:
            sig = b"\x89PNG\r\n\x1a\n"

            def _chunk(ctype: bytes, data: bytes) -> bytes:
                c = ctype + data
                return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

            ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
            raw = zlib.compress(b"\x00\xff\x00\x00")
            return sig + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", raw) + _chunk(b"IEND", b"")

        img1 = tmp_path / "img1.png"
        img2 = tmp_path / "img2.png"
        img1.write_bytes(_make_png())
        img2.write_bytes(_make_png())

        await handle(
            action="create",
            path="img_repl.pptx",
            slides=[{"title": "Slide"}],
        )
        await handle(
            action="insert_image",
            path="img_repl.pptx",
            slide_index=1,
            image_path=str(img1),
            left=1,
            top=1,
            width=2,
            height=2,
        )

        # Find the picture shape index
        from pptx import Presentation as _Presentation

        prs = _Presentation(str(tmp_path / "img_repl.pptx"))
        pic_idx = None
        for i, shape in enumerate(list(prs.slides)[0].shapes, 1):
            if shape.shape_type == 13:
                pic_idx = i
                break
        assert pic_idx is not None

        result = await handle(
            action="edit",
            path="img_repl.pptx",
            image_replacements=[
                {"slide_index": 1, "shape_index": pic_idx, "image_path": str(img2)},
            ],
        )
        assert "error" not in result
        assert result["images_replaced"] == 1

    @pytest.mark.asyncio
    async def test_image_replacement_non_picture_shape(self, tmp_path):
        await handle(
            action="create",
            path="img_np.pptx",
            slides=[{"title": "Text shape"}],
        )
        result = await handle(
            action="edit",
            path="img_np.pptx",
            image_replacements=[
                {"slide_index": 1, "shape_index": 1, "image_path": "img.png"},
            ],
        )
        assert "error" not in result
        assert result["images_replaced"] == 0

    @pytest.mark.asyncio
    async def test_image_replacement_shape_out_of_bounds(self, tmp_path):
        await handle(
            action="create",
            path="img_oob.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(
            action="edit",
            path="img_oob.pptx",
            image_replacements=[
                {"slide_index": 1, "shape_index": 999, "image_path": "img.png"},
            ],
        )
        assert "error" not in result
        assert result["images_replaced"] == 0


@_needs_pptx
class TestEditErrorMessageNewParams:
    @pytest.mark.asyncio
    async def test_error_message_includes_new_params(self, tmp_path):
        await handle(
            action="create",
            path="err_new.pptx",
            slides=[{"title": "Slide"}],
        )
        result = await handle(action="edit", path="err_new.pptx")
        assert "error" in result
        assert "template_fill" in result["error"]
        assert "table_format" in result["error"]
        assert "paragraph_edits" in result["error"]
        assert "placeholder_edits" in result["error"]
        assert "image_replacements" in result["error"]


@_needs_pptx
class TestEditOpsCap:
    @pytest.mark.asyncio
    async def test_too_many_table_edits(self, tmp_path):
        await handle(action="create", path="cap.pptx", slides=[{"title": "S"}])
        edits = [{"slide_index": 1, "table_index": 1, "row": 0, "col": 0, "value": "x"}] * (_MAX_EDIT_OPS + 1)
        result = await handle(action="edit", path="cap.pptx", table_edits=edits)
        assert "error" in result
        assert "Too many" in result["error"]

    @pytest.mark.asyncio
    async def test_too_many_paragraph_edits(self, tmp_path):
        await handle(action="create", path="cap2.pptx", slides=[{"title": "S"}])
        edits = [{"slide_index": 1, "shape_index": 1, "paragraphs": [{"text": "x"}]}] * (_MAX_EDIT_OPS + 1)
        result = await handle(action="edit", path="cap2.pptx", paragraph_edits=edits)
        assert "error" in result
        assert "Too many" in result["error"]


@_needs_pptx
class TestNewFeaturesDefinition:
    def test_template_fill_in_definition(self):
        props = DEFINITION["parameters"]["properties"]
        assert "template_fill" in props

    def test_table_format_in_definition(self):
        props = DEFINITION["parameters"]["properties"]
        assert "table_format" in props

    def test_paragraph_edits_in_definition(self):
        props = DEFINITION["parameters"]["properties"]
        assert "paragraph_edits" in props

    def test_placeholder_edits_in_definition(self):
        props = DEFINITION["parameters"]["properties"]
        assert "placeholder_edits" in props

    def test_image_replacements_in_definition(self):
        props = DEFINITION["parameters"]["properties"]
        assert "image_replacements" in props


# ---------------------------------------------------------------------------
# COM-only actions: verify they return the COM-only error on lib backend
# ---------------------------------------------------------------------------


@_needs_pptx
class TestTransitionsComOnly:
    @pytest.mark.asyncio
    async def test_transitions_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="t.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="transitions",
            path="t.pptx",
            slide_index=1,
            transition={"effect": "fade"},
        )
        assert "error" in result
        assert "COM backend" in result["error"]


@_needs_pptx
class TestAnimationsComOnly:
    @pytest.mark.asyncio
    async def test_animations_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="a.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="animations",
            path="a.pptx",
            slide_index=1,
            shape_index=1,
            effect_id=1,
        )
        assert "error" in result
        assert "COM backend" in result["error"]


@_needs_pptx
class TestEmbedChartComOnly:
    @pytest.mark.asyncio
    async def test_embed_chart_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="c.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="embed_chart",
            path="c.pptx",
            slide_index=1,
            chart_type="bar",
        )
        assert "error" in result
        assert "COM backend" in result["error"]


@_needs_pptx
class TestExportPdfComOnly:
    @pytest.mark.asyncio
    async def test_export_pdf_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="e.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="export_pdf",
            path="e.pptx",
        )
        assert "error" in result
        assert "COM backend" in result["error"]


@_needs_pptx
class TestSectionsComOnly:
    @pytest.mark.asyncio
    async def test_sections_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="s.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="sections",
            path="s.pptx",
            operation="list",
        )
        assert "error" in result
        assert "COM backend" in result["error"]


@_needs_pptx
class TestGroupShapesComOnly:
    @pytest.mark.asyncio
    async def test_group_shapes_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="g.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="group_shapes",
            path="g.pptx",
            slide_index=1,
            shape_indices=[1, 2],
        )
        assert "error" in result
        assert "COM backend" in result["error"]


@_needs_pptx
class TestAudioVideoComOnly:
    @pytest.mark.asyncio
    async def test_audio_video_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="av.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="audio_video",
            path="av.pptx",
            slide_index=1,
            media_path="clip.mp4",
        )
        assert "error" in result
        assert "COM backend" in result["error"]


@_needs_pptx
class TestSmartartComOnly:
    @pytest.mark.asyncio
    async def test_smartart_returns_com_only_error(self, tmp_path):
        await handle(action="create", path="sa.pptx", slides=[{"title": "S1"}])
        result = await handle(
            action="smartart",
            path="sa.pptx",
            slide_index=1,
            smartart_layout="basic_block_list",
            smartart_items=["Item 1", "Item 2"],
        )
        assert "error" in result
        assert "COM backend" in result["error"]


class TestComDispatchErrorHandling:
    @pytest.mark.asyncio
    async def test_dispatch_com_returns_error_dict_on_exception(self):
        """When a COM handler raises, _dispatch_com catches and returns error dict with details."""
        from unittest.mock import AsyncMock, MagicMock

        mock_manager = MagicMock()
        mock_manager.run_com = AsyncMock(side_effect=RuntimeError("Access denied by security policy"))
        mock_com_mod = MagicMock()
        mock_com_mod.get_manager.return_value = mock_manager
        mock_com_mod.COM_AVAILABLE = True

        with (
            patch("anteroom.tools.office_pptx.AVAILABLE", True),
            patch("anteroom.tools.office_pptx._BACKEND", "com"),
            patch("anteroom.tools.office_pptx._com_mod", mock_com_mod),
        ):
            result = await handle(action="edit", path="test.pptx", replacements=[{"old": "a", "new": "b"}])

        assert "error" in result
        assert "Access denied by security policy" in result["error"]
        assert "RuntimeError" in result["error"]

    @pytest.mark.asyncio
    async def test_open_pres_com_includes_exception_message(self):
        """_open_pres_com error includes exc type and message, not just type name."""
        from unittest.mock import MagicMock

        from anteroom.tools.office_pptx import _open_pres_com

        mock_manager = MagicMock()
        mock_ppt = MagicMock()
        mock_ppt.Presentations.Open.side_effect = Exception("File is locked by another process")
        mock_manager.get_app.return_value = mock_ppt

        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(b"dummy")
            path = f.name
        try:
            _, _, err = _open_pres_com(mock_manager, path, "test.pptx")
            assert err is not None
            assert "File is locked by another process" in err
            assert "Exception" in err
        finally:
            import os

            os.unlink(path)
